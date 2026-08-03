"""Design system: warm editorial food-catalogue look.

Paper-cream page, white cards with soft corners, editorial serif headlines, and a
signature colour per supplier drawn from their own packaging. The chrome stays
disciplined so the product photography carries the appetite.
"""
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ------------------------------------------------------------------- palette ----
INK = RGBColor(0x14, 0x24, 0x1F)        # headings, near-black green
BODY_TX = RGBColor(0x44, 0x52, 0x4C)    # body copy
MUTED = RGBColor(0x8A, 0x93, 0x8D)      # labels, captions
PAPER = RGBColor(0xFA, 0xF7, 0xF2)      # page — warm cream, not clinical white
WHITE = RGBColor(0xFF, 0xFF, 0xFF)      # cards
RULE = RGBColor(0xE6, 0xE1, 0xD8)       # hairlines on paper
GOLD = RGBColor(0xC9, 0xA2, 0x4E)       # small ornaments

# Signature colour per supplier, sampled from their packaging. Each carries
# white text at >= 4.5:1, so it can be used as a filled band.
SUPPLIER_COLORS = {
    'singha': RGBColor(0xB3, 0x62, 0x14),      # burnt amber
    'thainichi': RGBColor(0x1D, 0x5B, 0x63),   # deep teal
    'tmk': RGBColor(0x2F, 0x7A, 0x3E),         # fresh green
    'zek': RGBColor(0xB2, 0x3A, 0x22),         # terracotta
}
DEFAULT_ACCENT = RGBColor(0x2F, 0x7A, 0x3E)

HEAD = 'Cambria'
BODY = 'Calibri'

SW, SH = 10.0, 5.625
M = 0.40
TOP_RULE = 0.07
CONTENT_TOP = 1.22
CONTENT_BOTTOM = 5.16
HAIRLINE = 0.01
CARD_R = 0.035                          # soft corners — food, not audit


def mix(color, other, amount):
    """Blend `amount` of `color` into `other`."""
    return RGBColor(*(round(c * amount + o * (1 - amount))
                      for c, o in zip(color, other)))


def tint(color, amount=0.09):
    """Very pale wash of a supplier colour, for photo areas and panels."""
    return mix(color, WHITE, amount)


def deepen(color, amount=0.72):
    """Darker version of a supplier colour, for text on paper."""
    return mix(color, INK, amount)


def fmt_usd(v):
    return ('$%.3f' % v).replace('.', ',')


def fmt_uah(v):
    return ('%.2f ₴' % v).replace('.', ',')


# ---------------------------------------------------------------- primitives ----
def rect(slide, x, y, w, h, fill=None, radius=None, line=None, lw=0.75):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        sh.adjustments[0] = radius
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.text_frame.text = ''
    return sh


def floating_card(slide, path, cx, cy, w, h, angle=0, pad=0.10, shadow=None):
    """A small product photo on a white mat, tilted like a dropped print —
    the flat-lay device for showing several products at once without a grid.

    Frame and picture share the same centre, so python-pptx's own per-shape
    rotation (which pivots on that centre) keeps them visually locked together
    without needing an actual shape group.
    """
    if shadow:
        rect(slide, cx - w / 2 + 0.035, cy - h / 2 + 0.05, w, h, fill=shadow,
             radius=0.05).rotation = angle
    frame = rect(slide, cx - w / 2, cy - h / 2, w, h, fill=WHITE, radius=0.05)
    frame.rotation = angle
    pic = picture(slide, path, cx - w / 2 + pad, cy - h / 2 + pad, w - 2 * pad, h - 2 * pad)
    pic.rotation = angle
    return frame


def ring(slide, cx, cy, d, color, lw=1.4):
    """A stroked circle, no fill — for a quiet medallion motif on a colour field."""
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2),
                                Inches(d), Inches(d))
    sh.fill.background()
    sh.shadow.inherit = False
    sh.line.color.rgb = color
    sh.line.width = Pt(lw)
    sh.text_frame.text = ''
    return sh


def hline(slide, x, y, w, color=RULE):
    return rect(slide, x, y, w, HAIRLINE, fill=color)


def text(slide, x, y, w, h, runs, size=10, font=BODY, color=BODY_TX, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, italic=False,
         wrap=True, spc=0):
    """runs: a string, or a list of (text, {overrides}). `spc` is letter spacing in pt."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = runs if isinstance(runs, list) else [(runs, {})]
    # a "\n" inside a run is a soft break that ignores paragraph spacing — split
    # it into real paragraphs so line_spacing applies to every line
    paras = []
    for item in items:
        content, over = item if isinstance(item, tuple) else (item, {})
        paras += [(line, over) for line in content.split('\n')]
    for i, (content, over) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get('align', align)
        p.line_spacing = over.get('line_spacing', line_spacing)
        r = p.add_run()
        r.text = content
        f = r.font
        f.name = over.get('font', font)
        f.size = Pt(over.get('size', size))
        f.bold = over.get('bold', bold)
        f.italic = over.get('italic', italic)
        f.color.rgb = over.get('color', color)
        tracking = over.get('spc', spc)
        if tracking:
            r.font._rPr.set('spc', str(int(tracking * 100)))
    return box


def label(slide, x, y, w, content, color=MUTED, size=6.5, align=PP_ALIGN.LEFT,
          h=0.16):
    """Small letter-spaced caps — the workhorse label of a consulting layout.

    Keep word_wrap on: a wrap="none" body auto-sizes and LibreOffice then centres
    it on the box regardless of paragraph alignment.
    """
    return text(slide, x, y, w, h, content.upper(), size=size, color=color, bold=True,
                align=align, spc=0.9)


MIN_DPI = 90       # never render a picture coarser than this


def fit_size(path, bw, bh):
    """Size a picture would take in the box — lets a layout react to small shots."""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih, 1.0 / MIN_DPI)
    return iw * scale, ih * scale


def picture_cover(slide, path, bx, by, bw, bh, focus=0.5):
    """Fill the box edge to edge, cropping the overflow — for full-bleed hero art.

    python-pptx has no crop-on-insert, so the picture is placed oversized and the
    shape's own crop fields trim it back to the box.
    """
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = max(bw / iw, bh / ih)
    w, h = iw * scale, ih * scale
    pic = slide.shapes.add_picture(path, Inches(bx), Inches(by), Inches(w), Inches(h))
    over_w, over_h = (w - bw) / w, (h - bh) / h
    pic.crop_left = over_w * focus
    pic.crop_right = over_w * (1 - focus)
    pic.crop_top = over_h * focus
    pic.crop_bottom = over_h * (1 - focus)
    pic.left, pic.top = Inches(bx), Inches(by)
    pic.width, pic.height = Inches(bw), Inches(bh)
    return pic


def picture(slide, path, bx, by, bw, bh):
    """Fit (contain) inside the box, centred, but never blow a shot up past MIN_DPI.

    Some supplier shots are 43x64 thumbnails. Stretched to fill a box they turn
    to mush; held to their own resolution they stay crisp and simply sit smaller
    on the card, which reads as deliberate against a white photo area.
    """
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih, 1.0 / MIN_DPI)
    w, h = iw * scale, ih * scale
    return slide.shapes.add_picture(path, Inches(bx + (bw - w) / 2),
                                    Inches(by + (bh - h) / 2), Inches(w), Inches(h))


# ----------------------------------------------------------------- furniture ----
def blank(prs, accent=DEFAULT_ACCENT, bg=PAPER, rule=True):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, SW, SH, fill=bg)
    if rule:
        rect(s, 0, 0, SW, TOP_RULE, fill=accent)
    return s


def header(slide, title, eyebrow=None, tag=None, accent=DEFAULT_ACCENT):
    if eyebrow:
        label(slide, M, 0.34, 5.6, eyebrow, color=deepen(accent, 0.85), size=7.5)
    text(slide, M, 0.54, SW - 2 * M - 2.6, 0.46, title, size=20, font=HEAD, color=INK,
         bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if tag:
        text(slide, SW - M - 2.6, 0.54, 2.6, 0.46, tag, size=9, color=MUTED,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    hline(slide, M, 1.08, SW - 2 * M)


def footer(slide, page, note=None):
    hline(slide, M, 5.26, SW - 2 * M)
    if note:
        text(slide, M, 5.34, SW - 2 * M - 0.6, 0.16, note, size=6.5, color=MUTED,
             italic=True)
    text(slide, SW - M - 0.6, 5.34, 0.6, 0.16, str(page), size=7, color=MUTED,
         bold=True, align=PP_ALIGN.RIGHT)
