"""PANTRY LIGHT vi — the magazine opener.

Changes against v5, all of them requested:
  1. Pack shots are no longer cropped to a circle. Each keeps its own
     proportions inside a blush tile, with air around it.
  2. The Mishima contact line is gone.
  3. Euro is gone. The calculation carries $ and UAH — nothing else.
  4. Cost per gram is gone. Nothing is derived that the source file did not
     already state.
  5. Mishima and Takaokaya cards are now structurally identical to ZEK's:
     the same freight-scenario block, with their single 17 m3 line where ZEK
     has four.
  6/7. Profile pages open on full-bleed photography with a floating cream
     panel over it — the feature-opener move the reference magazines use.
     Mishima's is the manufacturer's own plant; Takaokaya has no reachable
     photograph of theirs, so its slot holds the captioned pack visualisation
     in exactly the same frame.
  8. The zest: photography now bleeds off every edge of the opener instead of
     sitting in a column, the display serif runs large over it, and the page
     furniture (eyebrow rule, small-caps caption, hairline grid) is set the
     way a masthead sets it.
"""
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

SRC, DST = "deck.pptx", "v6.pptx"

LINEN  = RGBColor(0xFA, 0xF7, 0xF2)
BLUSH  = RGBColor(0xF1, 0xDD, 0xD3)
CHERRY = RGBColor(0xC3, 0x27, 0x33)
INK    = RGBColor(0x23, 0x1C, 0x18)
GREY   = RGBColor(0x6A, 0x5F, 0x56)
MUTED  = RGBColor(0x9A, 0x91, 0x8A)
RULE   = RGBColor(0xDF, 0xD4, 0xC9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

SERIF, SANS = "Cambria", "Calibri"
MARGIN = 0.62
PAGE_W, PAGE_H = 13.33, 7.5


# ── primitives ────────────────────────────────────────────────────────────
def spc(run, points):
    run.font._rPr.set("spc", str(int(points * 100)))


def rect(slide, x, y, w, h, fill, rounded=False, adj=0.5, line=None, lw=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    if rounded:
        shape.adjustments[0] = adj
    return shape


def hairline(slide, x, y, w, color=RULE):
    return rect(slide, x, y, w, 0.012, color)


def ring(slide, cx, cy, d, color=CHERRY, lw=1.0, dashed=True):
    """Empty specimen ring — a slot held open for a photograph not yet taken."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2),
                                   Inches(cy - d / 2), Inches(d), Inches(d))
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(lw)
    if dashed:
        shape.line.dash_style = MSO_LINE.DASH
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, body, size=10, bold=False, italic=False, color=INK,
         font=SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0, line=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = [body] if isinstance(body, str) else body
    for i, item in enumerate(items):
        s, over = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = over.get("align", align)
        if line or over.get("line"):
            para.line_spacing = over.get("line", line)
        run = para.add_run()
        run.text = s
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
        sp = over.get("spacing", spacing)
        if sp:
            spc(run, sp)
    return box


def smallcaps(slide, x, y, w, label, size=8.4, color=CHERRY, align=PP_ALIGN.LEFT):
    return text(slide, x, y, w, 0.20, label, size=size, bold=True, color=color,
                align=align, spacing=0.9)


def picture_cover(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    box_r, im_r = w / h, iw / ih
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if abs(box_r - im_r) > 1e-3:
        if im_r > box_r:
            keep = box_r / im_r
            pic.crop_left = pic.crop_right = (1 - keep) / 2
        else:
            keep = im_r / box_r
            pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    return pic


def picture_contain(slide, path, x, y, w, h):
    """Whole image, never cropped — the pack keeps its own proportions."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                    Inches(pw), Inches(ph))


def eyebrow(slide, x, y, label, color=CHERRY):
    rect(slide, x, y + 0.082, 0.28, 0.016, color)
    text(slide, x + 0.40, y, 6.5, 0.20, label, size=9.4, bold=True, color=color, spacing=1.0)


def stamp(slide, x, y, label, size=9.2):
    w = 0.52 + 0.088 * len(label)
    rect(slide, x, y, w, 0.34, None, rounded=True, adj=0.5, line=CHERRY, lw=1.0)
    text(slide, x, y, w, 0.34, label, size=size, bold=True, color=CHERRY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=0.6)
    return w


def note(slide, msg, y, x=MARGIN, w=None, h=0.48):
    w = w if w is not None else PAGE_W - 2 * MARGIN - 0.75
    rect(slide, x, y, 0.022, h, CHERRY)
    text(slide, x + 0.20, y, w - 0.20, h, msg, size=8.1, color=GREY, line=1.18,
         anchor=MSO_ANCHOR.MIDDLE)


def page_no(slide, n, color=MUTED):
    text(slide, PAGE_W - MARGIN - 0.6, PAGE_H - 0.46, 0.6, 0.26, str(n),
         size=9, color=color, align=PP_ALIGN.RIGHT, font=SERIF)


def blank(prs, layout):
    slide = prs.slides.add_slide(layout)
    rect(slide, 0, 0, PAGE_W, PAGE_H, LINEN)
    return slide


def num(v, dec=2):
    return f"{v:,.{dec}f}".replace(",", " ").replace(".", ",")


# ── the opener: full-bleed photograph, floating cream panel ───────────────
def profile_slide(prs, layout, *, page, eyebrow_label, name_lines, subtitle,
                  fob, about_lines, stats, hero, hero_caption, tail=None, warning=None):
    slide = blank(prs, layout)
    picture_cover(slide, hero, 0, 0, PAGE_W, PAGE_H)

    PX, PY, PW = 0.78, 0.86, 6.42
    pad = 0.52
    inner = PW - 2 * pad

    # panel height is measured from its contents so the card never runs long
    h = pad
    h += 0.32                                   # eyebrow
    h += 0.60 * len(name_lines) + 0.06          # display name
    h += 0.36                                   # italic dek
    h += 0.52                                   # fob stamp
    h += 0.15 + 0.24                            # rule + label
    h += 0.198 * len(about_lines) + 0.14        # about
    h += 0.20                                   # closing rule
    h += 0.86                                   # stats band
    if tail:
        h += 0.52
    if warning:
        h += 0.62
    h += pad - 0.14
    rect(slide, PX, PY, PW, h, LINEN)

    x, y = PX + pad, PY + pad
    eyebrow(slide, x, y, eyebrow_label)
    y += 0.32

    text(slide, x, y, inner, 0.60 * len(name_lines),
         [(ln, {"size": 32, "bold": True, "color": INK, "font": SERIF, "line": 1.02})
          for ln in name_lines])
    y += 0.60 * len(name_lines) + 0.06

    text(slide, x, y, inner, 0.30, subtitle, size=13, italic=True, color=CHERRY, font=SERIF)
    y += 0.36

    stamp(slide, x, y, fob)
    y += 0.52

    hairline(slide, x, y, inner)
    y += 0.15
    smallcaps(slide, x, y, inner, "ПРО КОМПАНІЮ")
    y += 0.24
    text(slide, x, y, inner, 0.198 * len(about_lines) + 0.06,
         [(ln, {"size": 10.4, "color": INK, "font": SERIF, "line": 1.15})
          for ln in about_lines])
    y += 0.198 * len(about_lines) + 0.14
    hairline(slide, x, y, inner)
    y += 0.20

    sw = inner / len(stats)
    for i, (val, cap) in enumerate(stats):
        sx = x + i * sw
        if i:
            rect(slide, sx - 0.02, y + 0.02, 0.011, 0.58, RULE)
        p = 0.20 if i else 0
        text(slide, sx + p, y, sw - p - 0.10, 0.40, val, size=25, bold=True,
             color=CHERRY, font=SERIF)
        text(slide, sx + p, y + 0.40, sw - p - 0.10, 0.34, cap.split("\n"),
             size=8.0, color=GREY, line=1.12)
    y += 0.86

    if tail:
        label, value = tail
        smallcaps(slide, x, y, inner, label)
        text(slide, x, y + 0.23, inner, 0.26, value, size=10.2, color=INK, font=SERIF)
        y += 0.52

    if warning:
        note(slide, warning, y + 0.06, x=x, w=inner, h=0.50)

    # caption and folio ride a slim cream tab so they never fight the photograph
    tab_w, tab_h = 5.20, 0.36
    tab_x, tab_y = PAGE_W - MARGIN - tab_w, PAGE_H - MARGIN - tab_h + 0.16
    rect(slide, tab_x, tab_y, tab_w, tab_h, LINEN)
    text(slide, tab_x + 0.22, tab_y, tab_w - 0.90, tab_h, hero_caption,
         size=7.4, italic=True, color=GREY, font=SERIF,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, tab_x + tab_w - 0.60, tab_y, 0.42, tab_h, str(page),
         size=9, color=MUTED, font=SERIF, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


# ── cost page: one component, ZEK's four scenarios or the others' one ─────
def sku(name, grams, per_carton, desc, scenarios, photo=None, photo_note=None):
    """scenarios: [(label, $, UAH)] — first one gets the display treatment."""
    return dict(name=name, grams=grams, per_carton=per_carton, desc=desc,
                scenarios=scenarios, photo=photo, photo_note=photo_note)


CARD_TOP = 1.34
ROW = 0.285                              # height of one freight-scenario row


def draw_card(slide, item, x, w):
    T = CARD_TOP
    # One component, two tempos. ZEK quotes four freight scenarios; the other
    # two quote one. The rows ZEK spends on scenarios, they spend on the
    # photograph — so every card still ends on the same baseline.
    tile_h = 1.95 + ROW * (4 - len(item["scenarios"]))
    grow = tile_h - 1.95
    rect(slide, x, T, w, tile_h, BLUSH)
    if item["photo"]:
        picture_contain(slide, item["photo"], x + 0.20, T + 0.16, w - 0.40, tile_h - 0.32)
        if item["photo_note"]:
            text(slide, x, T + tile_h + 0.05, w, 0.18, item["photo_note"], size=6.7,
                 italic=True, color=GREY, font=SERIF, align=PP_ALIGN.CENTER)
    else:
        d = min(w - 1.10, tile_h - 1.10)
        ring(slide, x + w / 2, T + tile_h / 2 - 0.16, d)
        rect(slide, x + w / 2 - 0.045, T + tile_h / 2 - 0.205, 0.09, 0.09, CHERRY)
        text(slide, x, T + tile_h / 2 + d / 2 + 0.06, w, 0.24, "ОЧІКУЄМО ФОТО",
             size=8.0, bold=True, color=CHERRY, align=PP_ALIGN.CENTER, spacing=1.4)

    y = T + 2.10 + grow
    rect(slide, x, y + 0.055, 0.075, 0.075, CHERRY)
    text(slide, x + 0.16, y, w - 0.16, 0.18,
         f"{num(item['grams'], 0)} Г · {item['per_carton']} ШТ/КАРТ.",
         size=7.2, bold=True, color=INK, spacing=0.8)

    y += 0.24
    per_line = max(12, int(w / 0.092))
    lines = 1 + (len(item["name"]) - 1) // per_line
    text(slide, x, y, w, 0.28 * lines + 0.05, item["name"],
         size=13.5, bold=True, color=INK, font=SERIF, line=1.04)

    y += 0.28 * lines + 0.08
    text(slide, x, y, w, 0.56, item["desc"], size=8.2, color=GREY, line=1.18)

    y = T + 3.30 + grow
    hairline(slide, x, y, w)
    y += 0.12
    smallcaps(slide, x, y, w, f"СОБІВАРТІСТЬ · 1 ПАКЕТ {num(item['grams'], 0)} Г", size=7.4)

    head, *rest = item["scenarios"]
    label, usd, uah = head
    y += 0.22
    text(slide, x, y, w, 0.18, label, size=7.0, bold=True, color=MUTED, spacing=0.7)
    y += 0.18
    text(slide, x, y, w * 0.56, 0.44, usd, size=25, bold=True, color=CHERRY, font=SERIF)
    text(slide, x + w * 0.42, y + 0.13, w * 0.58, 0.28, uah, size=13, bold=True,
         color=INK, font=SERIF, align=PP_ALIGN.RIGHT)
    y += 0.54

    for label, usd, uah in rest:
        hairline(slide, x, y, w)
        text(slide, x, y + 0.08, w * 0.50, 0.20, label, size=7.4, color=GREY)
        text(slide, x + w * 0.40, y + 0.07, w * 0.26, 0.22, usd, size=9.0,
             color=INK, font=SERIF, align=PP_ALIGN.RIGHT)
        text(slide, x + w * 0.62, y + 0.07, w * 0.38, 0.22, uah, size=9.0, bold=True,
             color=INK, font=SERIF, align=PP_ALIGN.RIGHT)
        y += 0.285
    hairline(slide, x, y, w)
    return y


def card_slide(prs, layout, *, page, eyebrow_label, title, tag, items, warning):
    slide = blank(prs, layout)
    eyebrow(slide, MARGIN, 0.44, eyebrow_label)
    text(slide, MARGIN, 0.70, 8.2, 0.58, title, size=30, bold=True, color=INK, font=SERIF)

    tw = 0.52 + 0.086 * len(tag)
    rect(slide, PAGE_W - MARGIN - tw, 0.64, tw, 0.36, None, rounded=True, adj=0.5,
         line=CHERRY, lw=1.0)
    text(slide, PAGE_W - MARGIN - tw, 0.64, tw, 0.36, tag, size=8.8, bold=True,
         color=CHERRY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=0.5)

    n = len(items)
    gap = 0.34 if n == 3 else 0.26
    w = (PAGE_W - 2 * MARGIN - gap * (n - 1)) / n
    for i, item in enumerate(items):
        draw_card(slide, item, MARGIN + i * (w + gap), w)

    note(slide, warning, 6.76)
    page_no(slide, page)
    return slide


# ══ content ═══════════════════════════════════════════════════════════════
P = "assets/packs"

ZEK_FREIGHT_35 = [("40′ КОНТЕЙНЕР", "$0,7", "31,54 ₴"),
                  ("Збірний 34 м³", "$0,7", "33,68 ₴"),
                  ("20′ контейнер", "$0,9", "38,66 ₴"),
                  ("Збірний 17 м³", "$0,9", "40,95 ₴")]
ZEK_FREIGHT_70 = [("40′ КОНТЕЙНЕР", "$1,4", "63,86 ₴"),
                  ("Збірний 34 м³", "$1,5", "68,19 ₴"),
                  ("20′ контейнер", "$1,7", "78,29 ₴"),
                  ("Збірний 17 м³", "$1,8", "82,92 ₴")]

ZEK_35 = [
    sku("Seaweed Topping Chicken Floss", 35, 24,
        "Подрібнена норі з курячою стружкою — присипка до рису та супів. Збагачена DHA.",
        ZEK_FREIGHT_35, f"{P}/zek_chicken35.png"),
    sku("Seaweed Topping Vegetables", 35, 24,
        "Норі-присипка з овочами. Паковання орієнтоване на дитячу аудиторію.",
        ZEK_FREIGHT_35, f"{P}/zek_veg35.png"),
    sku("Seaweed Topping Sesame", 35, 24,
        "Норі-присипка з кунжутом — базовий смак напрямку, без м’ясних добавок.",
        ZEK_FREIGHT_35, f"{P}/zek_sesame35.png"),
]
ZEK_70 = [
    sku("Seaweed Topping Vegetables", 70, 24,
        "Овочева присипка у великому пакеті 70 г. Розрахована на родину або HoReCa.",
        ZEK_FREIGHT_70, f"{P}/zek_veg70.png"),
    sku("Seaweed Topping Sesame", 70, 24,
        "Кунжутна присипка у пакеті 70 г — удвічі більший обсяг за той самий смак.",
        ZEK_FREIGHT_70, f"{P}/zek_sesame70.png"),
    sku("Seaweed Topping Chicken", 70, 24,
        "Присипка з курячою стружкою, 70 г. Найдорожча позиція в розрахунку ZEK.",
        ZEK_FREIGHT_70, f"{P}/zek_chicken70.png"),
]


def one(name, grams, carton, usd, uah, desc, photo=None, pnote=None):
    """Mishima / Takaokaya: ZEK's block with their single priced scenario."""
    return sku(name, grams, carton, desc,
               [("ЗБІРНИЙ ВАНТАЖ 17 М³", f"${num(usd)}", f"{num(uah)} ₴")], photo, pnote)


MISHIMA_25 = [
    one("Furikake з креветкою", 25, 20, 0.72, 36.49,
        "Видима креветка та дрібна риба — джерело мікроелементів. 3–4 г на 200 г рису.",
        f"{P}/ms_shrimp.png"),
    one("Furikake васабі", 25, 20, 0.89, 45.10,
        "Освіжаючий аромат васабі. Підходить до смаженого, м’яса на грилі та сашимі.",
        f"{P}/ms_wasabi.png", "фото — коробка 80 г з каталогу"),
    one("Furikake зі смаком кімчі", 25, 20, 0.655, 33.19,
        "Гострий профіль кімчі. У роздрібному каталозі постачальника позиція відсутня."),
    one("Furikake з норі", 25, 20, 0.555, 28.12,
        "Кунжут і норі — базовий смак серії. Універсальна присипка до рису.",
        f"{P}/ms_nori.png"),
]
MISHIMA_50 = [
    one("Furikake з креветкою", 50, 64, 1.3797, 69.91,
        "Та сама рецептура у подвійній фасовці — 64 пакети в коробі.",
        f"{P}/ms_shrimp.png"),
    one("Furikake васабі", 50, 64, 1.6328, 82.74,
        "Найдорожча позиція Mishima в розрахунку собівартості.",
        f"{P}/ms_wasabi.png", "фото — коробка 80 г з каталогу"),
    one("Furikake зі смаком кімчі", 50, 64, 1.1969, 60.65,
        "Найдешевша з 50-грамових позицій у перерахунку на пакет."),
    one("Furikake з норі", 50, 64, 0.9953, 50.44,
        "Найдешевша позиція розрахунку Mishima.", f"{P}/ms_nori.png"),
]
TAKAOKAYA_BOTTLES = [
    one("Nori Furikake Bottle", 50, 10, 1.278, 86.02,
        "Класична норі-присипка у скляній пляшці fresh-lock."),
    one("Nori Wasabi Furikake Bottle", 50, 10, 1.278, 86.02,
        "Норі з гранулами японського васабі (регіон Адзуміно). Ціна — як у поз. 1.",
        f"{P}/tk_wasabi.png", "фото роздрібної версії 70 г"),
    one("Yuzu Kosho Furikake Bottle", 50, 10, 1.344, 90.46,
        "Юзу-кошьо з Кюсю: цитрусова цедра та зелений перець."),
]
TAKAOKAYA_45 = [
    one("Nori Katsuo Furikake", 45, 80, 1.416, 95.31,
        "Норі та стружка боніто — умамі-профіль."),
    one("Nori Tamago Furikake", 60, 80, 1.416, 95.31,
        "Норі та яєчний порошок. Найбільша фасовка серії за тією ж ціною."),
    one("Curry Furikake", 45, 80, 1.062, 71.48,
        "Каррі-присипка японського профілю спецій."),
    one("Garlic Chili Oil Furikake", 45, 80, 1.062, 71.48,
        "Часник і чилі-олія раю — гострий напрямок серії. Ціна — як у поз. 6."),
]


# ══ assemble ══════════════════════════════════════════════════════════════
prs = Presentation(SRC)
layout = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "Пустий слайд")

profile_slide(
    prs, layout, page=12, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 04",
    name_lines=["HanJin (Shanghai)", "Food Co., Ltd."],
    subtitle="Темпура, сендвічі та присипка з норі · Retail",
    fob="FOB QINGDAO",
    about_lines=[
        "ZEK — «Zesty Especial Kingdom», власний бренд Hanjin Food.",
        "Компанія заснована 2002 року у Вейхаї, головний офіс із 2014-го —",
        "у Шанхаї. Партнери групи — тайська Singha Group і корейська",
        "SPC Group.",
    ],
    stats=[("2002", "рік заснування\nу Вейхаї, Шаньдун"),
           ("№1", "частка ринку\nприправленої норі в КНР"),
           ("16", "країн\nекспорту")],
    hero="assets/hero/zek_still.jpg",
    hero_caption="лінійка Seaweed Topping · пакування постачальника",
    tail=("СИРОВИНА", "Норі з острова Чеджу (Корея) та плантацій під Вейхаєм"),
)

card_slide(prs, layout, page=13, eyebrow_label="ZEK · SEAWEED TOPPING · 1/2",
           title="Присипка 35 г — собівартість", tag="ЧОТИРИ СЦЕНАРІЇ ДОСТАВКИ",
           items=ZEK_35,
           warning="Ціни однакові для всіх трьох смаків у межах ваги — постачальник котирує "
                   "напрямок, а не артикул.")

card_slide(prs, layout, page=14, eyebrow_label="ZEK · SEAWEED TOPPING · 2/2",
           title="Присипка 70 г — собівартість", tag="ЧОТИРИ СЦЕНАРІЇ ДОСТАВКИ",
           items=ZEK_70,
           warning="Подвоєння ваги з 35 до 70 г підвищує собівартість пакета приблизно вдвічі.")

profile_slide(
    prs, layout, page=15, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 05",
    name_lines=["Dalian Mishima Foods", "Co., Ltd."],
    subtitle="Фурікаке, соуси та приправи Pan-Asian · Retail / OEM",
    fob="FOB DALIAN, КИТАЙ",
    about_lines=[
        "Мішіма — виробник соусів, фурікаке, супових баз та юзу-продукції",
        "(м. Далянь, пров. Ляонін, КНР). Асортимент за каталогом: класичні",
        "соуси, пастові соуси, фурікаке в пакетах і коробках, супові бази,",
        "очадзуке та юзу-серія.",
    ],
    stats=[("8 SKU", "фурікаке в розрахунку\nсобівартості (25 і 50 г)"),
           ("17 м³", "єдиний прорахований\nсценарій доставки"),
           ("12 міс.", "термін придатності\n(за каталогом)")],
    hero="assets/hero/mishima_plant.jpg",
    hero_caption="виробничий майданчик Mishima, Далянь · фото компанії",
    warning="У файлі розрахунку базис доставки вказано «FOB Bangkok», хоча країна походження — "
            "Китай. Невідповідність потребує уточнення до фіксації умов.",
)

card_slide(prs, layout, page=16, eyebrow_label="MISHIMA · FURIKAKE · 1/2",
           title="Furikake 25 г — собівартість", tag="ЗБІРНИЙ ВАНТАЖ 17 М³",
           items=MISHIMA_25,
           warning="Прорахований лише збірний вантаж 17 м³ — на відміну від 4-рівневого "
                   "порівняння ZEK. Для зіставлення за обсягом контейнера потрібно "
                   "домоделювати сценарії 40′ / 20′ / 34 м³.")

card_slide(prs, layout, page=17, eyebrow_label="MISHIMA · FURIKAKE · 2/2",
           title="Furikake 50 г — собівартість", tag="ЗБІРНИЙ ВАНТАЖ 17 М³",
           items=MISHIMA_50,
           warning="Смак «кімчі» відсутній у роздрібному каталозі постачальника — SKU та фото "
                   "потребують підтвердження.")

profile_slide(
    prs, layout, page=18, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 06",
    name_lines=["Takaokaya Co., Ltd.", "(бренд Kinjirushi)"],
    subtitle="Бутильоване фурікаке та васабі-приправи · Японія",
    fob="FOB ЯПОНІЯ (ПОРТ УТОЧНЮЄТЬСЯ)",
    about_lines=[
        "Заснована 1929 р., спеціалізується на продукції з автентичного",
        "японського васабі (регіон Адзуміно). Бренд Kinjirushi. Дочірня",
        "Takaokaya USA працює з 1986 р. Постачається у скляних «fresh-lock»",
        "пляшечках 45–60 г — преміум «Made in Japan».",
    ],
    stats=[("1929", "рік заснування\nбренду"),
           ("7 SKU", "фурікаке в розрахунку\nсобівартості"),
           ("45–60 г", "фасовка,\nскляна пляшка")],
    hero="assets/hero/takaokaya_still.jpg",
    hero_caption="візуалізація фасовки fresh-lock · фото виробництва не надані",
    warning="Юридична особа й пряме контактне джерело постачальника відсутні — профіль "
            "складено за відкритими даними бренду, потребує верифікації перед контрактом.",
)

card_slide(prs, layout, page=19, eyebrow_label="TAKAOKAYA · FURIKAKE · 1/2",
           title="Fresh-lock пляшки 50 г", tag="ЗБІРНИЙ ВАНТАЖ 17 М³",
           items=TAKAOKAYA_BOTTLES,
           warning="Позиції 1 і 2 мають однакову закупівельну ціну незалежно від смаку — типово "
                   "для непроробленого прайсу; варто запросити поартикульну котирувку.")

card_slide(prs, layout, page=20, eyebrow_label="TAKAOKAYA · FURIKAKE · 2/2",
           title="Фасовка 45–60 г", tag="ЗБІРНИЙ ВАНТАЖ 17 М³",
           items=TAKAOKAYA_45,
           warning="Пари 4/5 і 6/7 мають ідентичну ціну незалежно від ваги (45 г vs 60 г).")

sld = prs.slides._sldIdLst
for el in list(sld)[:7]:
    prs.part.drop_rel(el.rId)
    sld.remove(el)

prs.save(DST)
print("saved", DST, "slides:", len(sld))
