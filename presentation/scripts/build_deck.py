"""Add product imagery to the Mishima and Takaokaya slides of the supplier deck.

Keeps the deck's existing design language: warm paper background, terracotta
accents, Cambria headings / Calibri body, and the white "photo print" card with
a terracotta offset shadow used on the ZEK profile slide.
"""
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

SRC, DST = "deck.pptx", "out.pptx"

INK       = RGBColor(0x2A, 0x22, 0x1E)
BRICK     = RGBColor(0xB2, 0x39, 0x24)
TERRA     = RGBColor(0xC9, 0x75, 0x64)   # offset card behind photos (as on slide 1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLUSH     = RGBColor(0xFD, 0xE4, 0xDD)
MUTED     = RGBColor(0x8A, 0x81, 0x7A)
GREY      = RGBColor(0x6E, 0x67, 0x60)

FLAVOUR_DOT = {
    "Nori":             RGBColor(0x1E, 0x5B, 0x45),
    "Nori Wasabi":      RGBColor(0x74, 0xA8, 0x3E),
    "Yuzu Kosho":       RGBColor(0xC9, 0x97, 0x1F),
    "Nori Katsuo":      RGBColor(0xA8, 0x45, 0x2C),
    "Nori Tamago":      RGBColor(0xE0, 0xAE, 0x4B),
    "Curry":            RGBColor(0xC8, 0x7A, 0x1E),
    "Garlic Chili Oil": RGBColor(0xB0, 0x24, 0x18),
}

PAD = 0.13          # white frame around a photo, inches


def ratio(path):
    im = Image.open(path)
    return im.width / im.height


def drop(slide, *names):
    for shp in list(slide.shapes):
        if shp.name in names:
            shp._element.getparent().remove(shp._element)


def rect(slide, x, y, w, h, fill, rounded=False, adj=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    if rounded:
        shape.adjustments[0] = adj if adj is not None else 0.10
    return shape


def oval(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, runs, size=10, bold=False, color=INK, font="Calibri",
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=0, italic=False):
    """runs: a string, or a list of (text, {overrides}) tuples, one paragraph each."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        line, over = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = over.get("align", align)
        if space:
            para.space_after = Pt(space)
        run = para.add_run()
        run.text = line
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
    return box


def photo_card(slide, path, x, y, card_w, pad=PAD, shadow=True, offset=(0.05, 0.06)):
    """White photo-print card with a terracotta offset, sized from the image ratio.

    Returns the card's bottom edge in inches.
    """
    pw = card_w - 2 * pad
    ph = pw / ratio(path)
    card_h = ph + 2 * pad
    if shadow:
        rect(slide, x + offset[0], y + offset[1], card_w, card_h, TERRA, rounded=True, adj=0.04)
    rect(slide, x, y, card_w, card_h, WHITE, rounded=True, adj=0.04)
    slide.shapes.add_picture(path, Inches(x + pad), Inches(y + pad), Inches(pw), Inches(ph))
    return y + card_h


prs = Presentation(SRC)
s4, s5, s6, s7 = prs.slides[3], prs.slides[4], prs.slides[5], prs.slides[6]

# ─────────────────────────────────────────────────────────────── slide 4: Mishima
drop(s4, "Rectangle 26", "Picture 27", "Rectangle 28", "Picture 29",
     "Rectangle 30", "Picture 31")

bottom = photo_card(s4, "assets/mishima_lineup.jpg", 9.00, 0.85, 3.50)      # ~3.20
small_w, gap = 1.72, 0.06
photo_card(s4, "assets/mishima_sauces.jpg", 9.00, bottom + 0.16, small_w)
photo_card(s4, "assets/mishima_ochazuke_small.jpg", 9.00 + small_w + gap, bottom + 0.16, small_w)
text(s4, 8.95, 6.74, 3.60, 0.20,
     "Фото: роздрібний каталог Mishima (Retail Product Catalog)",
     size=7.5, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────── slide 5: Mishima cost
drop(s5, "Picture 7", "Picture 9", "Picture 11",
     "TextBox 8", "TextBox 10", "TextBox 12")

legend = [("assets/pack_shrimp.png", 0.55, "Креветка"),
          ("assets/pack_wasabi.png", 1.30, "Васабі"),
          ("assets/pack_nori.png",   2.05, "Норі")]
for path, x, label in legend:
    rect(s5, x - 0.07, 1.18, 0.72, 0.72, WHITE, rounded=True, adj=0.08)
    s5.shapes.add_picture(path, Inches(x - 0.01), Inches(1.24), Inches(0.60), Inches(0.60))
    text(s5, x - 0.20, 1.94, 0.95, 0.22, label, size=8.5, color=GREY, align=PP_ALIGN.CENTER)

col_x, col_w = 9.85, 2.95
bottom = photo_card(s5, "assets/mishima_rice.jpg", col_x, 1.25, col_w)             # ~3.58
bottom = photo_card(s5, "assets/mishima_ochazuke.jpg", col_x, bottom + 0.22, col_w)
rect(s5, col_x, 5.95, col_w, 0.50, BLUSH, rounded=True, adj=0.22)
text(s5, col_x + 0.15, 6.05, col_w - 0.30, 0.32,
     [("Фурікаке 25 / 50 г", {"size": 9.5, "bold": True, "color": BRICK, "font": "Cambria"}),
      ("фото: роздрібний каталог Mishima", {"size": 7.5, "color": GREY})],
     align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────── slide 6: Takaokaya
drop(s6, *[f"Rounded Rectangle {n}" for n in (25, 28, 31, 34, 37, 40, 43)],
     *[f"Rectangle {n}" for n in (26, 29, 32, 35, 38, 41, 44)],
     *[f"TextBox {n}" for n in (27, 30, 33, 36, 39, 42, 45)])

bottom = photo_card(s6, "assets/takaokaya_hero.jpg", 9.00, 0.85, 3.50)      # ~3.26
text(s6, 9.00, bottom + 0.10, 3.50, 0.20, "FURIKAKE · 7 СМАКІВ",
     size=8.5, bold=True, color=BRICK)

chips = list(FLAVOUR_DOT)
chip_w, chip_h, gx, gy = 1.72, 0.38, 0.06, 0.06
top = bottom + 0.36
for i, name in enumerate(chips[:6]):
    cx = 9.00 + (i % 2) * (chip_w + gx)
    cy = top + (i // 2) * (chip_h + gy)
    rect(s6, cx, cy, chip_w, chip_h, WHITE, rounded=True, adj=0.22)
    oval(s6, cx + 0.13, cy + chip_h / 2 - 0.045, 0.09, FLAVOUR_DOT[name])
    text(s6, cx + 0.30, cy, chip_w - 0.38, chip_h, name, size=9, bold=True,
         color=INK, anchor=MSO_ANCHOR.MIDDLE)

last = chips[6]
cy = top + 3 * (chip_h + gy)
rect(s6, 9.00, cy, chip_w * 2 + gx, chip_h, WHITE, rounded=True, adj=0.22)
oval(s6, 9.13, cy + chip_h / 2 - 0.045, 0.09, FLAVOUR_DOT[last])
text(s6, 9.30, cy, 3.10, chip_h, last, size=9, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)

text(s6, 8.95, 6.68, 3.60, 0.20,
     "Візуалізація фасовки fresh-lock · фото постачальника відсутні",
     size=7.5, color=MUTED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────── slide 7: Takaokaya cost
bottom = photo_card(s7, "assets/takaokaya_bottle.jpg", 11.00, 1.45, 1.80)   # ~3.74
rect(s7, 11.00, bottom + 0.16, 1.80, 0.92, BLUSH, rounded=True, adj=0.12)
text(s7, 11.10, bottom + 0.28, 1.60, 0.70,
     [("FRESH-LOCK", {"size": 10.5, "bold": True, "color": BRICK, "font": "Cambria"}),
      ("скляна пляшка", {"size": 8, "color": GREY}),
      ("45–60 г", {"size": 8, "color": GREY})],
     align=PP_ALIGN.CENTER)
text(s7, 11.00, bottom + 1.16, 1.80, 0.30,
     "Візуалізація фасовки — фото постачальника відсутні",
     size=7, color=MUTED, align=PP_ALIGN.CENTER)

prs.save(DST)
print("saved", DST)
