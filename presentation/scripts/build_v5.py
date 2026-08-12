"""PANTRY LIGHT — one design language across all nine slides.

The deck had been running two visual systems: ZEK (slides 1-3) still on the
original cream-card layout, Mishima/Takaokaya on a separately redesigned one.
That split — not the colour — was the real reason the deck looked unresolved.

v5 rebuilds every slide from scratch in a single system:

  * one ground: warm linen paper, never interrupted by a dark field. The black
    price block is gone; cost now sits on the paper in large cherry serif
    numerals over hairline rules, which is both prettier and easier to scan
    than white-on-black ever was.
  * one product form: a perfect circle on a blush halo, for every supplier.
    Where no photograph exists the circle stays, drawn as an empty outline.
  * one pair of voices: Cambria for what is said (names, figures, the italic
    dek), Calibri small-caps for what is labelled (units, eyebrows, captions).
  * one rhythm: identical left margin, eyebrow height and rule grid on all
    nine pages. ZEK's four freight scenarios and the other two suppliers'
    single 17 m3 scenario are the same component at two tempos, not two designs.

Full data legibility on the analytical slides is a hard requirement, so every
figure that was on the previous version is still here, in the same reading
order, at equal or larger size.
"""
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

SRC, DST = "deck.pptx", "v5.pptx"

# ── palette: five values, nothing else ────────────────────────────────────
LINEN  = RGBColor(0xFA, 0xF7, 0xF2)   # the ground, on every page
BLUSH  = RGBColor(0xF1, 0xDD, 0xD3)   # halo behind product portraits
CHERRY = RGBColor(0xC3, 0x27, 0x33)   # the single accent
INK    = RGBColor(0x23, 0x1C, 0x18)   # text only — never a field
GREY   = RGBColor(0x6A, 0x5F, 0x56)
MUTED  = RGBColor(0x9A, 0x91, 0x8A)
RULE   = RGBColor(0xDF, 0xD4, 0xC9)

SERIF, SANS = "Cambria", "Calibri"
MARGIN = 0.62
PAGE_W, PAGE_H = 13.33, 7.5
COL_X = 8.55                      # photography always bleeds off the right flank
COL_W = PAGE_W - COL_X
TEXT_W = COL_X - MARGIN - 0.38

FLAVOUR_DOT = {
    "Nori":             RGBColor(0x1E, 0x5B, 0x45),
    "Nori Wasabi":      RGBColor(0x74, 0xA8, 0x3E),
    "Yuzu Kosho":       RGBColor(0xC9, 0x97, 0x1F),
    "Nori Katsuo":      RGBColor(0xA8, 0x45, 0x2C),
    "Nori Tamago":      RGBColor(0xE0, 0xAE, 0x4B),
    "Curry":            RGBColor(0xC8, 0x7A, 0x1E),
    "Garlic Chili Oil": RGBColor(0xB0, 0x24, 0x18),
}


# ── primitives ────────────────────────────────────────────────────────────
def spc(run, points):
    run.font._rPr.set("spc", str(int(points * 100)))


def rect(slide, x, y, w, h, fill, rounded=False, adj=0.5, line=None, lw=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    if rounded:
        shape.adjustments[0] = adj
    return shape


def oval(slide, x, y, d, fill=None, line=None, lw=1.1):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def hairline(slide, x, y, w, color=RULE):
    """The rule does the work a filled panel used to do."""
    return rect(slide, x, y, w, 0.012, color)


def text(slide, x, y, w, h, body, size=10, bold=False, italic=False, color=INK,
         font=SANS, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0, line=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = [body] if isinstance(body, str) else body
    for i, item in enumerate(items):
        s, over = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = over.get("align", align)
        if line or over.get("line"):
            para.line_spacing = over.get("line", line)
        run = para.add_run()
        run.text = s
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
        sp = over.get("spacing", spacing)
        if sp:
            spc(run, sp)
    return box


def smallcaps(slide, x, y, w, label, size=8.4, color=CHERRY, align=PP_ALIGN.LEFT):
    """Every label in the deck is set this way: small, spaced, quiet."""
    return text(slide, x, y, w, 0.20, label, size=size, bold=True, color=color,
                align=align, spacing=0.9)


def picture_cover(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    box_ratio, im_ratio = w / h, iw / ih
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if abs(box_ratio - im_ratio) > 1e-3:
        if im_ratio > box_ratio:
            keep = box_ratio / im_ratio
            pic.crop_left = pic.crop_right = (1 - keep) / 2
        else:
            keep = im_ratio / box_ratio
            pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    return pic


def portrait(slide, path, cx, cy, halo_d, fill_d=None):
    """A pack shot as still life: perfect circle, blush halo, constant fill."""
    fill_d = fill_d or halo_d * 0.90
    oval(slide, cx - halo_d / 2, cy - halo_d / 2, halo_d, BLUSH)
    slide.shapes.add_picture(path, Inches(cx - fill_d / 2), Inches(cy - fill_d / 2),
                             Inches(fill_d), Inches(fill_d))


def empty_portrait(slide, cx, cy, halo_d, fill_d=None):
    """No photograph yet — an open place setting, not a hole in the page."""
    fill_d = fill_d or halo_d * 0.90
    oval(slide, cx - halo_d / 2, cy - halo_d / 2, halo_d, BLUSH)
    oval(slide, cx - fill_d / 2, cy - fill_d / 2, fill_d, None, line=CHERRY, lw=1.0)
    bw, bh = 0.40, 0.26
    rect(slide, cx - bw / 2, cy - bh / 2, bw, bh, None, rounded=True, adj=0.22,
         line=CHERRY, lw=1.0)
    rect(slide, cx - 0.07, cy - bh / 2 - 0.055, 0.14, 0.055, None, rounded=True, adj=0.3,
         line=CHERRY, lw=1.0)
    oval(slide, cx - 0.065, cy - 0.065 + 0.015, 0.13, None, line=CHERRY, lw=1.0)


def eyebrow(slide, x, y, label):
    rect(slide, x, y + 0.082, 0.28, 0.016, CHERRY)
    text(slide, x + 0.40, y, 6.5, 0.20, label, size=9.4, bold=True, color=CHERRY, spacing=1.0)


def stamp(slide, x, y, label, size=9.2, pad=0.52):
    w = pad + 0.088 * len(label)
    rect(slide, x, y, w, 0.34, None, rounded=True, adj=0.5, line=CHERRY, lw=1.0)
    text(slide, x, y, w, 0.34, label, size=size, bold=True, color=CHERRY,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=0.6)
    return w


def note(slide, msg, y, x=MARGIN, w=None, h=0.52):
    """A caveat in its own quiet register — marked by the accent rule, not by alarm."""
    w = w if w is not None else PAGE_W - 2 * MARGIN
    rect(slide, x, y, 0.022, h, CHERRY)
    text(slide, x + 0.20, y, w - 0.20, h, msg, size=8.2, color=GREY, line=1.20,
         anchor=MSO_ANCHOR.MIDDLE)
    return y + h


def page_no(slide, n):
    text(slide, PAGE_W - MARGIN - 0.6, PAGE_H - 0.46, 0.6, 0.26, str(n),
         size=9, color=MUTED, align=PP_ALIGN.RIGHT, font=SERIF)


def blank(prs, layout):
    slide = prs.slides.add_slide(layout)
    rect(slide, 0, 0, PAGE_W, PAGE_H, LINEN)
    return slide


def num(v, dec=2):
    return f"{v:,.{dec}f}".replace(",", " ").replace(".", ",")


# ── profile page ──────────────────────────────────────────────────────────
def profile_slide(prs, layout, *, page, eyebrow_label, name_lines, subtitle, fob,
                  about_lines, stats, column, column_caption=None,
                  tail=None, flavours=None, warning=None):
    slide = blank(prs, layout)

    kind, payload = column
    if kind == "photo":
        picture_cover(slide, payload, COL_X, 0, COL_W, PAGE_H)
    elif kind == "image":
        picture_cover(slide, payload, COL_X, 0, COL_W, PAGE_H)
    elif kind == "circles":
        rect(slide, COL_X, 0, COL_W, PAGE_H, BLUSH)
        for path, cx, cy, d in payload:
            oval(slide, cx - d / 2, cy - d / 2, d, LINEN)
            slide.shapes.add_picture(path, Inches(cx - d * 0.90 / 2), Inches(cy - d * 0.90 / 2),
                                     Inches(d * 0.90), Inches(d * 0.90))
    if column_caption:
        # keep clear of the page number, which lives in the same bottom-right corner
        text(slide, COL_X + 0.22, PAGE_H - 0.50, COL_W - 0.44 - 0.66, 0.24, column_caption,
             size=6.8, italic=True, color=GREY, font=SERIF, align=PP_ALIGN.CENTER)

    y = 0.46
    eyebrow(slide, MARGIN, y, eyebrow_label)
    y += 0.34

    text(slide, MARGIN, y, TEXT_W, 0.62 * len(name_lines),
         [(ln, {"size": 33, "bold": True, "color": INK, "font": SERIF, "line": 1.02})
          for ln in name_lines])
    y += 0.62 * len(name_lines) + 0.06

    text(slide, MARGIN, y, TEXT_W, 0.30, subtitle, size=13, italic=True,
         color=CHERRY, font=SERIF)
    y += 0.44

    stamp(slide, MARGIN, y, fob)
    y += 0.56

    hairline(slide, MARGIN, y, TEXT_W)
    y += 0.15
    smallcaps(slide, MARGIN, y, TEXT_W, "ПРО КОМПАНІЮ")
    y += 0.26
    about_h = 0.203 * len(about_lines)
    text(slide, MARGIN, y, TEXT_W, about_h + 0.06,
         [(ln, {"size": 10.6, "color": INK, "font": SERIF, "line": 1.16})
          for ln in about_lines])
    y += about_h + 0.11
    hairline(slide, MARGIN, y, TEXT_W)
    y += 0.22

    sw = TEXT_W / len(stats)
    for i, (val, cap) in enumerate(stats):
        sx = MARGIN + i * sw
        if i:
            rect(slide, sx - 0.02, y + 0.02, 0.011, 0.60, RULE)
        pad = 0.22 if i else 0
        text(slide, sx + pad, y, sw - pad - 0.12, 0.40, val, size=26, bold=True,
             color=CHERRY, font=SERIF)
        text(slide, sx + pad, y + 0.42, sw - pad - 0.12, 0.34, cap.split("\n"),
             size=8.2, color=GREY, line=1.14)
    y += 0.82

    if flavours:
        y += 0.14
        cols = 4
        gap = 0.08
        cw = (TEXT_W - gap * (cols - 1)) / cols
        ch, rg = 0.29, 0.05
        for i, nm in enumerate(flavours):
            c, r = i % cols, i // cols
            cx = MARGIN + c * (cw + gap)
            cy = y + r * (ch + rg)
            oval(slide, cx + 0.015, cy + ch / 2 - 0.042, 0.084, FLAVOUR_DOT[nm])
            text(slide, cx + 0.18, cy, cw - 0.19, ch, nm, size=8.2, bold=True,
                 color=INK, anchor=MSO_ANCHOR.MIDDLE)
        y += (-(-len(flavours) // cols)) * (ch + rg)

    if tail:
        y += 0.24
        label, value = tail
        smallcaps(slide, MARGIN, y, TEXT_W, label)
        text(slide, MARGIN, y + 0.25, TEXT_W, 0.26, value, size=10.4, color=INK, font=SERIF)
        y += 0.60
        hairline(slide, MARGIN, y, TEXT_W)   # the column closes on the same grid
        y += 0.18

    if warning:
        wy = min(max(y + 0.20, 6.28), PAGE_H - 0.30 - 0.52)
        note(slide, warning, wy, x=MARGIN, w=TEXT_W)
    page_no(slide, page)
    return slide


# ── cost page ─────────────────────────────────────────────────────────────
def sku(name, grams, per_carton, desc, headline, rows, photo=None, photo_note=None):
    """headline: (scenario label, $ value, UAH value). rows: [(label, mid, right)]
    — ZEK passes three further freight scenarios, Mishima/Takaokaya pass their
    two derived figures. One component, two tempos."""
    return dict(name=name, grams=grams, per_carton=per_carton, desc=desc,
                headline=headline, rows=rows, photo=photo, photo_note=photo_note)


CARD_TOP = 1.40


def draw_card(slide, item, x, w):
    T = CARD_TOP
    halo = 1.34
    cx = x + w / 2
    cy = T + 0.70

    if item["photo"]:
        portrait(slide, item["photo"], cx, cy, halo)
    else:
        empty_portrait(slide, cx, cy, halo)

    cap_y = cy + halo / 2 + 0.09
    if item["photo"] and item["photo_note"]:
        text(slide, x, cap_y, w, 0.18, item["photo_note"], size=6.7, italic=True,
             color=GREY, font=SERIF, align=PP_ALIGN.CENTER)
    elif not item["photo"]:
        text(slide, x, cap_y, w, 0.18, "ОЧІКУЄМО ФОТО", size=7.4, bold=True,
             color=CHERRY, align=PP_ALIGN.CENTER, spacing=0.9)

    y = T + 1.72
    oval(slide, x, y + 0.032, 0.082, CHERRY)
    text(slide, x + 0.16, y, w - 0.16, 0.18,
         f"{num(item['grams'], 0)} Г · {item['per_carton']} ШТ/КАРТ.",
         size=7.2, bold=True, color=INK, spacing=0.8)

    y += 0.26
    per_line = max(12, int(w / 0.092))
    name_h = 0.275 * (1 + (len(item["name"]) - 1) // per_line)
    text(slide, x, y, w, name_h + 0.05, item["name"],
         size=13.5, bold=True, color=INK, font=SERIF, line=1.04)

    y = T + 2.36
    text(slide, x, y, w, 0.60, item["desc"], size=8.2, color=GREY, line=1.18)

    # ── the figures: on paper, in cherry serif, over rules ──
    y = T + 3.06
    hairline(slide, x, y, w)
    y += 0.13
    smallcaps(slide, x, y, w, f"СОБІВАРТІСТЬ · 1 ПАКЕТ {num(item['grams'], 0)} Г", size=7.4)

    y += 0.24
    label, usd, uah = item["headline"]
    text(slide, x, y, w, 0.18, label, size=7.0, bold=True, color=MUTED, spacing=0.7)
    y += 0.20
    text(slide, x, y, w * 0.56, 0.44, usd, size=25, bold=True, color=CHERRY, font=SERIF)
    text(slide, x + w * 0.42, y + 0.13, w * 0.58, 0.28, uah,
         size=13, bold=True, color=INK, font=SERIF, align=PP_ALIGN.RIGHT)

    y += 0.52
    for row in item["rows"]:
        hairline(slide, x, y, w)
        label, mid, right = row
        text(slide, x, y + 0.08, w * 0.50, 0.20, label, size=7.4, color=GREY)
        if mid:
            text(slide, x + w * 0.42, y + 0.07, w * 0.26, 0.22, mid,
                 size=9.0, color=INK, font=SERIF, align=PP_ALIGN.RIGHT)
        text(slide, x + w * 0.62, y + 0.07, w * 0.38, 0.22, right,
             size=9.0, bold=True, color=INK, font=SERIF, align=PP_ALIGN.RIGHT)
        y += 0.285
    hairline(slide, x, y, w)
    return y


def card_slide(prs, layout, *, page, eyebrow_label, title, tag, items, warning):
    slide = blank(prs, layout)
    eyebrow(slide, MARGIN, 0.46, eyebrow_label)
    text(slide, MARGIN, 0.72, 8.2, 0.60, title, size=30, bold=True, color=INK, font=SERIF)

    tw = 0.52 + 0.086 * len(tag)
    rect(slide, PAGE_W - MARGIN - tw, 0.66, tw, 0.36, None, rounded=True, adj=0.5,
         line=CHERRY, lw=1.0)
    text(slide, PAGE_W - MARGIN - tw, 0.66, tw, 0.36, tag, size=8.8, bold=True,
         color=CHERRY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=0.5)

    n = len(items)
    gap = 0.34 if n == 3 else 0.26
    w = (PAGE_W - 2 * MARGIN - gap * (n - 1)) / n
    bottom = CARD_TOP
    for i, item in enumerate(items):
        bottom = draw_card(slide, item, MARGIN + i * (w + gap), w)

    note(slide, warning, min(max(bottom + 0.26, 6.42), PAGE_H - 0.30 - 0.52),
         x=MARGIN, w=PAGE_W - 2 * MARGIN)
    page_no(slide, page)
    return slide


# ══ content ═══════════════════════════════════════════════════════════════
C = "assets/circles"

ZEK_35 = [
    sku("Seaweed Topping Chicken Floss", 35, 24,
        "Подрібнена норі з курячою стружкою — присипка до рису та супів. Збагачена DHA.",
        ("40′ КОНТЕЙНЕР", "$0,7", "31,54 ₴"),
        [("Збірний 34 м³", "$0,7", "33,68 ₴"),
         ("20′ контейнер", "$0,9", "38,66 ₴"),
         ("Збірний 17 м³", "$0,9", "40,95 ₴")],
        f"{C}/zek_chicken35.png"),
    sku("Seaweed Topping Vegetables", 35, 24,
        "Норі-присипка з овочами. Паковання орієнтоване на дитячу аудиторію.",
        ("40′ КОНТЕЙНЕР", "$0,7", "31,54 ₴"),
        [("Збірний 34 м³", "$0,7", "33,68 ₴"),
         ("20′ контейнер", "$0,9", "38,66 ₴"),
         ("Збірний 17 м³", "$0,9", "40,95 ₴")],
        f"{C}/zek_veg35.png"),
    sku("Seaweed Topping Sesame", 35, 24,
        "Норі-присипка з кунжутом — базовий смак напрямку, без м’ясних добавок.",
        ("40′ КОНТЕЙНЕР", "$0,7", "31,54 ₴"),
        [("Збірний 34 м³", "$0,7", "33,68 ₴"),
         ("20′ контейнер", "$0,9", "38,66 ₴"),
         ("Збірний 17 м³", "$0,9", "40,95 ₴")],
        f"{C}/zek_sesame35.png"),
]

ZEK_70 = [
    sku("Seaweed Topping Vegetables", 70, 24,
        "Овочева присипка у великому пакеті 70 г. Розрахована на родину або HoReCa.",
        ("40′ КОНТЕЙНЕР", "$1,4", "63,86 ₴"),
        [("Збірний 34 м³", "$1,5", "68,19 ₴"),
         ("20′ контейнер", "$1,7", "78,29 ₴"),
         ("Збірний 17 м³", "$1,8", "82,92 ₴")],
        f"{C}/zek_veg70.png"),
    sku("Seaweed Topping Sesame", 70, 24,
        "Кунжутна присипка у пакеті 70 г — удвічі більший обсяг за той самий смак.",
        ("40′ КОНТЕЙНЕР", "$1,4", "63,86 ₴"),
        [("Збірний 34 м³", "$1,5", "68,19 ₴"),
         ("20′ контейнер", "$1,7", "78,29 ₴"),
         ("Збірний 17 м³", "$1,8", "82,92 ₴")],
        f"{C}/zek_sesame70.png"),
    sku("Seaweed Topping Chicken", 70, 24,
        "Присипка з курячою стружкою, 70 г. Найдорожча позиція в розрахунку ZEK.",
        ("40′ КОНТЕЙНЕР", "$1,4", "63,86 ₴"),
        [("Збірний 34 м³", "$1,5", "68,19 ₴"),
         ("20′ контейнер", "$1,7", "78,29 ₴"),
         ("Збірний 17 м³", "$1,8", "82,92 ₴")],
        f"{C}/zek_chicken70.png"),
]


def ms(name, grams, carton, usd, eur, uah, desc, photo=None, pnote=None):
    return sku(name, grams, carton, desc,
               ("ЗБІРНИЙ ВАНТАЖ 17 М³", f"${num(usd)}", f"{num(uah)} ₴"),
               [("С/С, €/од.", "", f"€{num(eur)}"),
                ("Питома собівартість", "", f"{num(uah / grams)} ₴/г")],
               photo, pnote)


MISHIMA_25 = [
    ms("Furikake з креветкою", 25, 20, 0.72, 0.8108, 36.49,
       "Видима креветка та дрібна риба — джерело мікроелементів. 3–4 г на 200 г рису.",
       f"{C}/ms_shrimp.png"),
    ms("Furikake васабі", 25, 20, 0.89, 1.0022, 45.10,
       "Освіжаючий аромат васабі. Підходить до смаженого, м’яса на грилі та сашимі.",
       f"{C}/ms_wasabi.png", "фото — коробка 80 г з каталогу"),
    ms("Furikake зі смаком кімчі", 25, 20, 0.655, 0.7376, 33.19,
       "Гострий профіль кімчі. У роздрібному каталозі постачальника позиція відсутня."),
    ms("Furikake з норі", 25, 20, 0.555, 0.6250, 28.12,
       "Кунжут і норі — базовий смак серії. Універсальна присипка до рису.",
       f"{C}/ms_nori.png"),
]

MISHIMA_50 = [
    ms("Furikake з креветкою", 50, 64, 1.3797, 1.5537, 69.91,
       "Та сама рецептура у подвійній фасовці — 64 пакети в коробі.",
       f"{C}/ms_shrimp.png"),
    ms("Furikake васабі", 50, 64, 1.6328, 1.8387, 82.74,
       "Найдорожча позиція Mishima в розрахунку: ₴1,65 за грам.",
       f"{C}/ms_wasabi.png", "фото — коробка 80 г з каталогу"),
    ms("Furikake зі смаком кімчі", 50, 64, 1.1969, 1.3478, 60.65,
       "Найнижча питома собівартість серед 50-грамових позицій."),
    ms("Furikake з норі", 50, 64, 0.9953, 1.1208, 50.44,
       "Найдешевша позиція розрахунку — ₴1,01 за грам.",
       f"{C}/ms_nori.png"),
]

TAKAOKAYA_BOTTLES = [
    ms("Nori Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
       "Класична норі-присипка у скляній пляшці fresh-lock."),
    ms("Nori Wasabi Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
       "Норі з гранулами японського васабі (регіон Адзуміно). Ціна — як у поз. 1.",
       f"{C}/tk_wasabi.png", "фото роздрібної версії 70 г"),
    ms("Yuzu Kosho Furikake Bottle", 50, 10, 1.344, 2.0103, 90.46,
       "Юзу-кошьо з Кюсю: цитрусова цедра та зелений перець."),
]

TAKAOKAYA_45 = [
    ms("Nori Katsuo Furikake", 45, 80, 1.416, 2.1179, 95.31,
       "Норі та стружка боніто — умамі-профіль. Найвища питома с/с: ₴2,12/г."),
    ms("Nori Tamago Furikake", 60, 80, 1.416, 2.1179, 95.31,
       "Норі та яєчний порошок. Найбільша фасовка серії за тією ж ціною."),
    ms("Curry Furikake", 45, 80, 1.062, 1.5884, 71.48,
       "Каррі-присипка японського профілю спецій."),
    ms("Garlic Chili Oil Furikake", 45, 80, 1.062, 1.5884, 71.48,
       "Часник і чилі-олія раю — гострий напрямок серії. Ціна — як у поз. 6."),
]


# ══ assemble ══════════════════════════════════════════════════════════════
prs = Presentation(SRC)
layout = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "Пустий слайд")

profile_slide(
    prs, layout, page=12, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 04",
    name_lines=["HanJin (Shanghai)", "Food Co., Ltd."],
    subtitle="Темпура, сендвічі та присипка з норі · Retail",
    fob="FOB QINGDAO",
    about_lines=[
        "ZEK — «Zesty Especial Kingdom», власний бренд Hanjin Food.",
        "Компанія заснована 2002 року у Вейхаї, головний офіс із 2014-го —",
        "у Шанхаї. Партнери групи — тайська Singha Group і корейська",
        "SPC Group.",
    ],
    stats=[("2002", "рік заснування\nу Вейхаї, Шаньдун"),
           ("№1", "частка ринку\nприправленої норі в КНР"),
           ("16", "країн\nекспорту")],
    column=("circles", [
        (f"{C}/zek_nori.png",    10.30, 2.10, 2.30),
        (f"{C}/zek_veg35.png",   11.93, 4.43, 1.75),
        (f"{C}/zek_tempura.png",  9.73, 5.28, 1.45),
    ]),
    tail=("СИРОВИНА", "Норі з острова Чеджу (Корея) та плантацій під Вейхаєм"),
)

card_slide(prs, layout, page=13, eyebrow_label="ZEK · SEAWEED TOPPING · 1/2",
           title="Присипка 35 г — собівартість",
           tag="ЧОТИРИ СЦЕНАРІЇ ДОСТАВКИ", items=ZEK_35,
           warning="Ціни однакові для всіх трьох смаків у межах ваги — постачальник котирує "
                   "напрямок, а не артикул. Розрахунок наведено для чотирьох сценаріїв "
                   "завантаження, тож ZEK можна зіставляти за обсягом контейнера.")

card_slide(prs, layout, page=14, eyebrow_label="ZEK · SEAWEED TOPPING · 2/2",
           title="Присипка 70 г — собівартість",
           tag="ЧОТИРИ СЦЕНАРІЇ ДОСТАВКИ", items=ZEK_70,
           warning="Подвоєння ваги з 35 до 70 г підвищує собівартість пакета приблизно вдвічі — "
                   "питома вартість за грам залишається на рівні 35-грамової фасовки.")

profile_slide(
    prs, layout, page=15, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 05",
    name_lines=["Dalian Mishima Foods", "Co., Ltd."],
    subtitle="Фурікаке, соуси та приправи Pan-Asian · Retail / OEM",
    fob="FOB DALIAN, КИТАЙ",
    about_lines=[
        "Мішіма — виробник соусів, фурікаке, супових баз та юзу-продукції",
        "(м. Далянь, пров. Ляонін, КНР). Асортимент за каталогом: класичні",
        "соуси, пастові соуси, фурікаке в пакетах і коробках, супові бази,",
        "очадзуке, юзу-серія, фреш-лок пляшки.",
        "Контакт: foods@mishima.com.cn · +86 411 87611161",
    ],
    stats=[("8 SKU", "фурікаке в розрахунку\nсобівартості (25 і 50 г)"),
           ("17 м³", "єдиний прорахований\nсценарій доставки"),
           ("12 міс.", "термін придатності\n(за каталогом)")],
    column=("photo", "assets/mishima_profile_hero.jpg"),
    warning="У файлі розрахунку базис доставки вказано «FOB Bangkok», хоча країна походження — "
            "Китай. Невідповідність потребує уточнення у постачальника до фіксації умов.",
)

card_slide(prs, layout, page=16, eyebrow_label="MISHIMA · FURIKAKE · 1/2",
           title="Furikake 25 г — собівартість",
           tag="ЗБІРНИЙ ВАНТАЖ 17 М³", items=MISHIMA_25,
           warning="Розрахунок побудований лише для збірного вантажу 17 м³ — на відміну від "
                   "4-рівневого порівняння ZEK. Для зіставлення постачальників за обсягом "
                   "контейнера потрібно домоделювати сценарії 40′ / 20′ / 34 м³.")

card_slide(prs, layout, page=17, eyebrow_label="MISHIMA · FURIKAKE · 2/2",
           title="Furikake 50 г — собівартість",
           tag="ЗБІРНИЙ ВАНТАЖ 17 М³", items=MISHIMA_50,
           warning="Смак «кімчі» відсутній у роздрібному каталозі постачальника — SKU та фото "
                   "потребують підтвердження. Фото 25 і 50 г ідентичні: постачальник показує "
                   "одну ілюстрацію пакування на смак.")

profile_slide(
    prs, layout, page=18, eyebrow_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 06",
    name_lines=["Takaokaya Co., Ltd.", "(бренд Kinjirushi)"],
    subtitle="Бутильоване фурікаке та васабі-приправи · Японія",
    fob="FOB ЯПОНІЯ (ПОРТ УТОЧНЮЄТЬСЯ)",
    about_lines=[
        "Заснована 1929 р., спеціалізується на продукції з автентичного",
        "японського васабі (регіон Адзуміно). Бренд Kinjirushi. Дочірня",
        "Takaokaya USA працює з 1986 р. Постачається у скляних «fresh-lock»",
        "пляшечках 45–60 г — преміум «Made in Japan» на відміну від",
        "пакетованого фурікаке з Китаю.",
    ],
    stats=[("1929", "рік заснування\nбренду"),
           ("7 SKU", "фурікаке в розрахунку\nсобівартості"),
           ("45–60 г", "фасовка,\nскляна пляшка")],
    column=("image", "assets/takaokaya_column.jpg"),
    column_caption="візуалізація фасовки fresh-lock · фото постачальника відсутні",
    flavours=list(FLAVOUR_DOT),
    warning="Юридична особа й пряме контактне джерело постачальника в матеріалах компанії "
            "відсутні — профіль складено за відкритими даними бренду, потребує верифікації "
            "перед контрактом.",
)

card_slide(prs, layout, page=19, eyebrow_label="TAKAOKAYA · FURIKAKE · 1/2",
           title="Fresh-lock пляшки 50 г",
           tag="ЗБІРНИЙ ВАНТАЖ 17 М³", items=TAKAOKAYA_BOTTLES,
           warning="Позиції 1 і 2 мають однакову закупівельну ціну незалежно від смаку — типово "
                   "для непроробленого прайсу; варто запросити поартикульну котирувку.")

card_slide(prs, layout, page=20, eyebrow_label="TAKAOKAYA · FURIKAKE · 2/2",
           title="Фасовка 45–60 г",
           tag="ЗБІРНИЙ ВАНТАЖ 17 М³", items=TAKAOKAYA_45,
           warning="Пари 4/5 і 6/7 мають ідентичну ціну незалежно від ваги (45 г vs 60 г). Питома "
                   "собівартість Takaokaya (₴1,59–2,12/г) у 1,4–2,1 раза вища за Mishima "
                   "(₴0,99–1,80/г) — очікувана премія «Made in Japan».")

# the nine new pages were appended after the seven originals; drop the originals
sld = prs.slides._sldIdLst
ids = list(sld)
for el in ids[:7]:
    prs.part.drop_rel(el.rId)
    sld.remove(el)

prs.save(DST)
print("saved", DST, "slides:", len(sld))
