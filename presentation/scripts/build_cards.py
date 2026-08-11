"""Replace the two cost tables with per-SKU cards (photo + 17 m³ self-cost).

Mirrors the ZEK card slides already in the deck: pack shot on a blush panel,
weight/carton badge, product name, short description, and a terracotta cost
block — here with the single logistics scenario these two suppliers have.
"""
import copy

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SRC, DST = "out.pptx", "final.pptx"

INK    = RGBColor(0x2A, 0x22, 0x1E)
BRICK  = RGBColor(0xB2, 0x39, 0x24)
BLUSH  = RGBColor(0xFD, 0xE4, 0xDD)
PANEL  = RGBColor(0xF9, 0xEF, 0xEA)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFA, 0xF7, 0xF2)
GREY   = RGBColor(0x6E, 0x67, 0x60)
MUTED  = RGBColor(0x8A, 0x81, 0x7A)
LINE   = RGBColor(0xE7, 0xDF, 0xD8)
ROSE   = RGBColor(0xC4, 0x8B, 0x7C)

MARGIN, TOP = 0.55, 1.28


def num(v, dec=2, suffix=""):
    s = f"{v:,.{dec}f}".replace(",", " ").replace(".", ",")
    return s + suffix


# ─────────────────────────────────────────────────────── SKU data (from Self_Cost.xlsx)
def sku(name, grams, per_carton, usd, eur, uah, desc, photo=None, photo_note=None):
    return dict(name=name, grams=grams, per_carton=per_carton, usd=usd, eur=eur,
                uah=uah, desc=desc, photo=photo, photo_note=photo_note)


MS_SHRIMP = "assets/sku_ms_shrimp.png"
MS_NORI = "assets/sku_ms_nori.png"
MS_WASABI = "assets/sku_ms_wasabi.png"
TK_WASABI = "assets/sku_tk_wasabi.png"

MISHIMA_25 = [
    sku("Furikake з креветкою", 25, 20, 0.72, 0.8108, 36.49,
        "Видима креветка та дрібна риба — джерело мікроелементів. 3–4 г на 200 г рису.",
        MS_SHRIMP),
    sku("Furikake васабі", 25, 20, 0.89, 1.0022, 45.10,
        "Освіжаючий аромат васабі. Підходить до смаженого, м’яса на грилі та сашимі.",
        MS_WASABI, "фото — коробка 80 г з каталогу"),
    sku("Furikake зі смаком кімчі", 25, 20, 0.655, 0.7376, 33.19,
        "Гострий профіль кімчі. У роздрібному каталозі постачальника позиція відсутня.",
        None),
    sku("Furikake з норі", 25, 20, 0.555, 0.6250, 28.12,
        "Кунжут і норі — базовий смак серії. Універсальна присипка до рису.",
        MS_NORI),
]

MISHIMA_50 = [
    sku("Furikake з креветкою", 50, 64, 1.3797, 1.5537, 69.91,
        "Та сама рецептура у подвійній фасовці — 64 пакети в коробі.", MS_SHRIMP),
    sku("Furikake васабі", 50, 64, 1.6328, 1.8387, 82.74,
        "Найдорожча позиція Mishima в розрахунку: ₴1,65 за грам.",
        MS_WASABI, "фото — коробка 80 г з каталогу"),
    sku("Furikake зі смаком кімчі", 50, 64, 1.1969, 1.3478, 60.65,
        "Найнижча питома собівартість серед 50-грамових позицій.", None),
    sku("Furikake з норі", 50, 64, 0.9953, 1.1208, 50.44,
        "Найдешевша позиція розрахунку — ₴1,01 за грам.", MS_NORI),
]

TAKAOKAYA_BOTTLES = [
    sku("Nori Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
        "Класична норі-присипка у скляній пляшці fresh-lock.", None),
    sku("Nori Wasabi Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
        "Норі з гранулами японського васабі (регіон Адзуміно). Ціна — як у поз. 1.",
        TK_WASABI, "фото роздрібної версії 70 г"),
    sku("Yuzu Kosho Furikake Bottle", 50, 10, 1.344, 2.0103, 90.46,
        "Юзу-кошьо з Кюсю: цитрусова цедра та зелений перець.", None),
]

TAKAOKAYA_45 = [
    sku("Nori Katsuo Furikake", 45, 80, 1.416, 2.1179, 95.31,
        "Норі та стружка боніто — умамі-профіль. Найвища питома с/с: ₴2,12/г.", None),
    sku("Nori Tamago Furikake", 60, 80, 1.416, 2.1179, 95.31,
        "Норі та яєчний порошок. Найбільша фасовка серії за тією ж ціною.", None),
    sku("Curry Furikake", 45, 80, 1.062, 1.5884, 71.48,
        "Каррі-присипка японського профілю спецій.", None),
    sku("Garlic Chili Oil Furikake", 45, 80, 1.062, 1.5884, 71.48,
        "Часник і чилі-олія раю — гострий напрямок серії. Ціна — як у поз. 6.", None),
]


# ─────────────────────────────────────────────────────────────────── drawing helpers
def spc(run, points):
    """Character spacing (not exposed by python-pptx)."""
    run.font._rPr.set("spc", str(int(points * 100)))


def rect(slide, x, y, w, h, fill, rounded=False, adj=0.06, border=None, bw=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border
        shape.line.width = Pt(bw)
    shape.shadow.inherit = False
    if rounded:
        shape.adjustments[0] = adj
    return shape


def oval(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, body, size=10, bold=False, color=INK, font="Calibri",
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0, line=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = [body] if isinstance(body, str) else body
    for i, item in enumerate(items):
        s, over = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = over.get("align", align)
        if line or over.get("line"):
            para.line_spacing = over.get("line", line)
        run = para.add_run()
        run.text = s
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        sp = over.get("spacing", spacing)
        if sp:
            spc(run, sp)
    return box


def picture_contain(slide, path, x, y, w, h):
    """Place an image inside a box without distorting it."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                             Inches(pw), Inches(ph))


# ─────────────────────────────────────────────────────────────────────── card render
def draw_card(slide, item, x, w):
    pad = 0.16
    y = TOP
    photo_h = 1.95
    card_bottom = 6.72

    rect(slide, x, y, w, card_bottom - y, WHITE, rounded=True, adj=0.035, border=LINE)

    # photo panel
    px, pw = x + 0.05, w - 0.10
    rect(slide, px, y + 0.05, pw, photo_h, PANEL)
    note_h = 0.20 if item["photo_note"] else 0.0
    if item["photo"]:
        picture_contain(slide, item["photo"], px + 0.14, y + 0.13,
                        pw - 0.28, photo_h - 0.16 - note_h)
    else:
        text(slide, px, y + 0.05 + photo_h / 2 - 0.14, pw, 0.28, "ОЧІКУЄМО ФОТО",
             size=9, bold=True, color=ROSE, align=PP_ALIGN.CENTER, spacing=1.2)
    if item["photo_note"]:
        text(slide, px, y + photo_h - 0.09, pw, 0.16, item["photo_note"],
             size=6.5, color=MUTED, align=PP_ALIGN.CENTER)

    cy = y + photo_h + 0.24
    inner_x, inner_w = x + pad, w - 2 * pad

    # badge
    oval(slide, inner_x, cy + 0.035, 0.10, BRICK)
    text(slide, inner_x + 0.18, cy, inner_w - 0.18, 0.18,
         f"{num(item['grams'], 0)} Г · {item['per_carton']} ШТ/КАРТ.",
         size=7.5, bold=True, color=INK, spacing=0.8)

    # name — the description follows it directly, so estimate how many lines it takes
    cy += 0.28
    per_line = max(12, int(inner_w / 0.097))          # Cambria bold 13 pt
    lines = 1 + (len(item["name"]) - 1) // per_line
    name_h = 0.26 * lines
    text(slide, inner_x, cy, inner_w, name_h + 0.04, item["name"],
         size=13, bold=True, color=INK, font="Cambria", line=1.05)

    # description
    cy += name_h + 0.10
    text(slide, inner_x, cy, inner_w, 0.72, item["desc"], size=8.5, color=GREY, line=1.12)

    # cost block
    label_y = 5.02
    text(slide, inner_x, label_y, inner_w, 0.18,
         f"СОБІВАРТІСТЬ · 1 ПАКЕТ {num(item['grams'], 0)} Г",
         size=7.5, bold=True, color=BRICK, spacing=0.6)

    block_y = label_y + 0.24
    rect(slide, inner_x, block_y, inner_w, 0.80, BRICK)
    text(slide, inner_x + 0.14, block_y + 0.09, inner_w - 0.28, 0.16,
         "ЗБІРНИЙ ВАНТАЖ 17 М³", size=6.5, bold=True, color=WHITE, spacing=0.7)
    text(slide, inner_x + 0.14, block_y + 0.28, inner_w * 0.55, 0.42,
         f"${num(item['usd'])}", size=20, bold=True, color=WHITE, font="Cambria")
    text(slide, inner_x + inner_w * 0.45, block_y + 0.40, inner_w * 0.55 - 0.14, 0.26,
         f"{num(item['uah'])} ₴", size=12.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # detail rows
    row_y = block_y + 0.92
    for label, value in (("С/С, €/од.", f"€{num(item['eur'])}"),
                         ("Питома собівартість", f"{num(item['uah'] / item['grams'])} ₴/г")):
        rect(slide, inner_x, row_y, inner_w, 0.012, LINE)
        text(slide, inner_x, row_y + 0.07, inner_w * 0.62, 0.20, label, size=7.5, color=GREY)
        text(slide, inner_x + inner_w * 0.42, row_y + 0.05, inner_w * 0.58, 0.22, value,
             size=9.5, bold=True, color=INK, align=PP_ALIGN.RIGHT)
        row_y += 0.30


def card_slide(prs, layout, kicker, title, subtitle, items, footnote):
    slide = prs.slides.add_slide(layout)
    rect(slide, 0, 0, 13.33, 7.5, PAPER)
    rect(slide, 0, 0, 13.33, 0.12, BRICK)

    text(slide, MARGIN, 0.32, 6.0, 0.28, kicker, size=10.5, bold=True, color=BRICK, spacing=0.6)
    text(slide, MARGIN, 0.58, 7.6, 0.55, title, size=26, bold=True, color=INK, font="Cambria")

    pill_w = 4.0
    rect(slide, 13.33 - MARGIN - pill_w, 0.55, pill_w, 0.50, BLUSH, rounded=True, adj=0.28)
    text(slide, 13.33 - MARGIN - pill_w + 0.15, 0.68, pill_w - 0.30, 0.30, subtitle,
         size=10, bold=True, color=BRICK, align=PP_ALIGN.CENTER)

    n = len(items)
    gap = 0.22 if n == 4 else 0.28
    w = (13.33 - 2 * MARGIN - gap * (n - 1)) / n
    for i, item in enumerate(items):
        draw_card(slide, item, MARGIN + i * (w + gap), w)

    if footnote:
        text(slide, MARGIN, 6.88, 12.23, 0.36, footnote, size=8, color=GREY, line=1.15)
    return slide


# ────────────────────────────────────────────────────────────────────────── assemble
prs = Presentation(SRC)
blank = next(l for l in prs.slide_masters[0].slide_layouts if not l.placeholders or True
             if l.name == "Пустий слайд")

sld_lst = prs.slides._sldIdLst

new = [
    card_slide(prs, blank, "MISHIMA · FURIKAKE · 1/2",
               "Furikake 25 г — собівартість",
               "Збірний вантаж 17 м³ · один сценарій", MISHIMA_25,
               "⚠ Розрахунок побудований лише для збірного вантажу 17 м³ — на відміну від 4-рівневого порівняння в профілі ZEK. "
               "Для зіставлення постачальників за обсягом контейнера потрібно домоделювати сценарії 40′ / 20′ / 34 м³."),
    card_slide(prs, blank, "MISHIMA · FURIKAKE · 2/2",
               "Furikake 50 г — собівартість",
               "Збірний вантаж 17 м³ · один сценарій", MISHIMA_50,
               "⚠ Смак «кімчі» відсутній у роздрібному каталозі постачальника — SKU та фото потребують підтвердження. "
               "Фото 25 і 50 г ідентичні: постачальник показує одну ілюстрацію пакування на смак."),
    card_slide(prs, blank, "TAKAOKAYA · FURIKAKE · 1/2",
               "Fresh-lock пляшки 50 г",
               "Збірний вантаж 17 м³ · один сценарій", TAKAOKAYA_BOTTLES,
               "⚠ Позиції 1 і 2 мають однакову закупівельну ціну незалежно від смаку — типово для непроробленого прайсу; "
               "варто запросити поартикульну котирувку."),
    card_slide(prs, blank, "TAKAOKAYA · FURIKAKE · 2/2",
               "Фасовка 45–60 г",
               "Збірний вантаж 17 м³ · один сценарій", TAKAOKAYA_45,
               "⚠ Пари 4/5 і 6/7 мають ідентичну ціну незалежно від ваги (45 г vs 60 г). Питома собівартість Takaokaya "
               "(₴1,59–2,12/г) у 1,4–2,1 раза вища за Mishima (₴0,99–1,80/г) — очікувана премія «Made in Japan»."),
]

# drop the two table slides now that the replacements exist (indexes 4 and 6)
ids = list(sld_lst)
for idx in (6, 4):
    el = ids[idx]
    prs.part.drop_rel(el.rId)
    sld_lst.remove(el)

# order: ZEK ×3, Mishima profile + its 2 card slides, Takaokaya profile + its 2
ids = list(sld_lst)
order = [ids[0], ids[1], ids[2],
         ids[3], ids[5], ids[6],
         ids[4], ids[7], ids[8]]
for el in order:
    sld_lst.remove(el)
for el in order:
    sld_lst.append(el)

# renumber the page footers sequentially (the deck starts at 12)
for i, slide in enumerate(prs.slides, start=12):
    footer = None
    for shp in slide.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip().isdigit() and shp.left > Inches(11.5):
            footer = shp
    if footer is not None:
        run = footer.text_frame.paragraphs[0].runs[0]
        run.text = str(i)
    else:
        text(slide, 12.73, 7.10, 0.50, 0.30, str(i), size=10, color=GREY, align=PP_ALIGN.RIGHT)

prs.save(DST)
print("saved", DST, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
