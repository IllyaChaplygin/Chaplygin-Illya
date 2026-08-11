"""Approximate renderer for visual QA (LibreOffice is unusable in this sandbox).

Draws shapes, pictures and wrapped text from a .pptx with PIL. Font metrics come
from Liberation (metric-compatible with Arial, slightly wider than Calibri), so
text that fits here fits in PowerPoint too.
"""
import sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU = 914400
DPI = 110

FONTS = {
    ("Calibri", False): "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ("Calibri", True): "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ("Cambria", False): "/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf",
    ("Cambria", True): "/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf",
}
_cache = {}


def font(name, bold, pt):
    key = (name if name in ("Cambria",) else "Calibri", bool(bold))
    px = max(6, round(pt * DPI / 72))
    ck = key + (px,)
    if ck not in _cache:
        _cache[ck] = ImageFont.truetype(FONTS[key], px)
    return _cache[ck]


def px(v):
    return v / EMU * DPI


def rgb(color, default=(0, 0, 0)):
    try:
        c = color.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return default


def shape_fill(shp):
    try:
        if shp.fill.type is not None and shp.fill.type == 1:
            return rgb(shp.fill.fore_color)
    except Exception:
        pass
    return None


def draw_shape(d, img, shp, ox=0, oy=0):
    try:
        x, y = px(shp.left) + ox, px(shp.top) + oy
        w, h = px(shp.width), px(shp.height)
    except TypeError:
        return
    st = str(shp.shape_type)

    if "PICTURE" in st:
        try:
            im = Image.open(__import__("io").BytesIO(shp.image.blob)).convert("RGB")
            im = im.resize((max(1, round(w)), max(1, round(h))), Image.LANCZOS)
            img.paste(im, (round(x), round(y)))
            d.rectangle([x, y, x + w, y + h], outline=(150, 150, 150))
        except Exception as e:
            d.rectangle([x, y, x + w, y + h], fill=(200, 200, 200))
        return

    if "LINE" in st or "CONNECTOR" in st:
        d.line([x, y, x + w, y + h], fill=(210, 205, 198), width=1)
        return

    fill = shape_fill(shp)
    if fill:
        name = (shp.name or "").lower()
        try:
            adj = shp.adjustments[0] if len(shp.adjustments) else 0
        except Exception:
            adj = 0
        if "oval" in name:
            d.ellipse([x, y, x + w, y + h], fill=fill)
        elif "rounded" in name:
            r = max(2, min(w, h) * (adj if adj else 0.1))
            d.rounded_rectangle([x, y, x + w, y + h], radius=r, fill=fill)
        else:
            d.rectangle([x, y, x + w, y + h], fill=fill)

    if not shp.has_text_frame:
        return

    tf = shp.text_frame
    paras = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        r0 = runs[0]
        pt = r0.font.size.pt if r0.font.size else 12
        paras.append({
            "text": "".join(r.text for r in runs),
            "pt": pt,
            "bold": bool(r0.font.bold),
            "name": r0.font.name or "Calibri",
            "color": rgb(r0.font.color, (30, 30, 30)),
            "align": p.alignment,
        })
    if not paras:
        return

    ml = px(tf.margin_left or 0)
    mr = px(tf.margin_right or 0)
    mt = px(tf.margin_top or 0)
    avail = max(10, w - ml - mr)

    lines = []
    for p in paras:
        f = font(p["name"], p["bold"], p["pt"])
        words = p["text"].split(" ")
        cur = ""
        wrapped = []
        for word in words:
            trial = word if not cur else cur + " " + word
            if d.textlength(trial, font=f) <= avail or not cur:
                cur = trial
            else:
                wrapped.append(cur)
                cur = word
        wrapped.append(cur)
        for ln in wrapped:
            lines.append((ln, f, p))

    lh = [f.size * 1.22 for _, f, _ in lines]
    total = sum(lh)
    anchor = tf.vertical_anchor
    if anchor == MSO_ANCHOR.MIDDLE:
        cy = y + (h - total) / 2
    elif anchor == MSO_ANCHOR.BOTTOM:
        cy = y + h - total
    else:
        cy = y + mt

    for (ln, f, p), step in zip(lines, lh):
        tw = d.textlength(ln, font=f)
        if p["align"] == PP_ALIGN.CENTER:
            tx = x + ml + (avail - tw) / 2
        elif p["align"] == PP_ALIGN.RIGHT:
            tx = x + w - mr - tw
        else:
            tx = x + ml
        d.text((tx, cy), ln, font=f, fill=p["color"])
        cy += step


def render(path, out_prefix, only=None):
    prs = Presentation(path)
    W, H = round(px(prs.slide_width)), round(px(prs.slide_height))
    made = []
    for i, slide in enumerate(prs.slides, 1):
        if only and i not in only:
            continue
        img = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(img)
        for shp in slide.shapes:
            if str(shp.shape_type).startswith("GROUP"):
                for sub in shp.shapes:
                    draw_shape(d, img, sub)
            else:
                draw_shape(d, img, shp)
        name = f"{out_prefix}-{i}.png"
        img.save(name)
        made.append(name)
    print("\n".join(made))


if __name__ == "__main__":
    only = {int(a) for a in sys.argv[3:]} or None
    render(sys.argv[1], sys.argv[2], only)
