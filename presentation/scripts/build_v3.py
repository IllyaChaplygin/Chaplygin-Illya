"""v3 redesign: glossy-magazine palette on top of the v2 full-bleed layout.

v2 borrowed the Ajinomoto reference's grammar (full-bleed photo columns, dark
panels, price-as-one-big-number) but stayed inside a muted terracotta/ink
range that read as flat and grey. v3 keeps the v2 layout math untouched and
repaints it: a more saturated editorial red, a true near-black instead of a
brown-black, a warm gold accent used only on dark grounds (labels, the
warning icon, the photo-frame line), a serif display face for headlines
(paired with the existing bold sans for kickers/labels — the classic
magazine masthead/caption contrast), and a black-vs-red split between
"photo in hand" and "photo pending" cards so four placeholders in a row read
as a deliberate two-tone system rather than a wall of identical boxes.

Built fresh from deck.pptx. Slides 1-3 (ZEK) are untouched.
"""
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

SRC, DST = "deck.pptx", "v3.pptx"

INK       = RGBColor(0x18, 0x12, 0x0E)   # true near-black, warm undertone
INK_SOFT  = RGBColor(0x2A, 0x21, 0x1B)   # placeholder panel — a shade off INK
BRICK     = RGBColor(0xC7, 0x2A, 0x1C)   # saturated editorial red (was muddier terracotta)
BRICK_D   = RGBColor(0x8A, 0x1C, 0x12)   # deep red for photo-credit strips
PAPER     = RGBColor(0xFA, 0xF7, 0xF2)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CREAM     = RGBColor(0xF3, 0xEA, 0xDF)
GOLD      = RGBColor(0xC9, 0xA1, 0x5A)   # accent — dark grounds only, never on paper
INKTXT    = RGBColor(0x18, 0x12, 0x0E)
GREY      = RGBColor(0x6A, 0x5F, 0x56)
MUTED     = RGBColor(0x9A, 0x91, 0x8A)
LINE      = RGBColor(0xE7, 0xDF, 0xD8)
LINE_ON_INK = RGBColor(0xFF, 0xFF, 0xFF)

MARGIN = 0.55
PAGE_W, PAGE_H = 13.33, 7.5

FLAVOUR_DOT = {
    "Nori":             RGBColor(0x1E, 0x5B, 0x45),
    "Nori Wasabi":      RGBColor(0x74, 0xA8, 0x3E),
    "Yuzu Kosho":       RGBColor(0xC9, 0x97, 0x1F),
    "Nori Katsuo":      RGBColor(0xA8, 0x45, 0x2C),
    "Nori Tamago":      RGBColor(0xE0, 0xAE, 0x4B),
    "Curry":            RGBColor(0xC8, 0x7A, 0x1E),
    "Garlic Chili Oil": RGBColor(0xB0, 0x24, 0x18),
}


# ───────────────────────────────────────────────────────────────── helpers
def spc(run, points):
    run.font._rPr.set("spc", str(int(points * 100)))


def rect(slide, x, y, w, h, fill, rounded=False, adj=0.05, line=None, lw=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    if rounded:
        shape.adjustments[0] = adj
    return shape


def oval(slide, x, y, d, fill=None, line=None, lw=1.25, dh=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(dh or d))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, body, size=10, bold=False, color=INKTXT, font="Calibri",
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=0, line=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    items = [body] if isinstance(body, str) else body
    for i, item in enumerate(items):
        s, over = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = over.get("align", align)
        if line or over.get("line"):
            para.line_spacing = over.get("line", line)
        run = para.add_run()
        run.text = s
        f = run.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        sp = over.get("spacing", spacing)
        if sp:
            spc(run, sp)
    return box


def picture_cover(slide, path, x, y, w, h):
    """Crop-to-fill a picture into an exact box (no letterboxing)."""
    iw, ih = Image.open(path).size
    box_ratio, im_ratio = w / h, iw / ih
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if abs(box_ratio - im_ratio) < 1e-3:
        return pic
    if im_ratio > box_ratio:            # image wider than box -> crop sides
        keep = box_ratio / im_ratio
        pic.crop_left = pic.crop_right = (1 - keep) / 2
    else:                                # image taller -> crop top/bottom
        keep = im_ratio / box_ratio
        pic.crop_top = pic.crop_bottom = (1 - keep) / 2
    return pic


def picture_contain(slide, path, x, y, w, h):
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(path, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                    Inches(pw), Inches(ph))


def kicker(slide, x, y, label, color=BRICK, dot=BRICK):
    # gold reads well only on dark grounds — on paper, brick stays the marker colour
    rect(slide, x, y + 0.035, 0.13, 0.13, dot)
    text(slide, x + 0.22, y, 6.5, 0.22, label, size=10.5, bold=True, color=color, spacing=0.7)


def page_no(slide, n, on_dark_bar=None):
    if on_dark_bar:
        x, y, w, h = on_dark_bar
        text(slide, x, y, w, h, str(n), size=9.5, color=CREAM, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    else:
        text(slide, PAGE_W - MARGIN - 0.5, PAGE_H - 0.40, 0.5, 0.3, str(n),
             size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def placeholder_photo(slide, x, y, w, h, note=None):
    """Camera glyph drawn from primitives — no font/emoji dependency.

    Sits on the INK_SOFT panel (see draw_card), so the accent reads as gold
    rather than the cream used elsewhere on brick — a deliberate cue that
    this card is a different state (pending), not a smaller version of the
    photo cards.
    """
    cx, cy = x + w / 2, y + h / 2
    bw, bh = 0.62, 0.42
    rect(slide, cx - bw / 2, cy - bh / 2 - 0.02, bw, bh, None, rounded=True, adj=0.18,
         line=GOLD, lw=1.5)
    rect(slide, cx - 0.11, cy - bh / 2 - 0.09, 0.22, 0.09, None, rounded=True, adj=0.3,
         line=GOLD, lw=1.5)
    oval(slide, cx - 0.11, cy - bh / 2 - 0.02 + 0.06, 0.22, line=GOLD, lw=1.5)
    text(slide, x, cy + bh / 2 + 0.10, w, 0.22, "ОЧІКУЄМО ФОТО",
         size=9.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER, spacing=1.0)
    if note:
        text(slide, x + 0.1, cy + bh / 2 + 0.34, w - 0.2, 0.3, note,
             size=7, color=RGBColor(0xA8, 0x9C, 0x8C), align=PP_ALIGN.CENTER, line=1.1)


def warn_icon(slide, x, y, size=0.26):
    """Small gold triangle + '!' — a drawn mark instead of the emoji glyph,
    which renders inconsistently across fonts and reads as a cartoon sticker
    next to serif headlines and a gold/black palette."""
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 Inches(x), Inches(y), Inches(size), Inches(size * 0.9))
    tri.fill.solid()
    tri.fill.fore_color.rgb = GOLD
    tri.line.fill.background()
    tri.shadow.inherit = False
    text(slide, x, y + size * 0.18, size, size * 0.6, "!", size=int(size * 46),
         bold=True, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x + size + 0.16


def warning_bar(slide, msg, y, h=0.5, x=MARGIN, w=None):
    w = w if w is not None else PAGE_W - 2 * MARGIN
    rect(slide, x, y, w, h, INK, rounded=True, adj=0.10)
    text_x = warn_icon(slide, x + 0.20, y + h / 2 - 0.15)
    text_w = (x + w - 0.75) - text_x
    text(slide, text_x, y, text_w, h, msg, size=8.3, color=CREAM, line=1.15,
         anchor=MSO_ANCHOR.MIDDLE)
    return y + h


def blank_slide(prs, layout):
    slide = prs.slides.add_slide(layout)
    rect(slide, 0, 0, PAGE_W, PAGE_H, PAPER)
    return slide


def num(v, dec=2):
    return f"{v:,.{dec}f}".replace(",", " ").replace(".", ",")


# ═══════════════════════════════════════════════════════════════ profile slide
def profile_slide(prs, layout, *, page, kicker_label, name_lines, subtitle, fob,
                   photo, photo_dark_bg, photo_caption, about_lines, stats,
                   flavours, warning):
    slide = blank_slide(prs, layout)

    photo_x = 8.75
    photo_w = PAGE_W - photo_x
    if photo_dark_bg:
        rect(slide, photo_x, 0, photo_w, PAGE_H, INK)
    picture_cover(slide, photo, photo_x, 0, photo_w, PAGE_H)
    if photo_caption:
        rect(slide, photo_x, PAGE_H - 0.42, photo_w, 0.42, INK)
        text(slide, photo_x + 0.18, PAGE_H - 0.42, photo_w - 0.36, 0.42, photo_caption,
             size=7, color=CREAM, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, line=1.05)

    col_w = photo_x - MARGIN - 0.35
    y = 0.42
    kicker(slide, MARGIN, y, kicker_label)
    y += 0.30

    name_h = 0.62 * len(name_lines)
    text(slide, MARGIN, y, col_w, name_h,
         [(ln, {"size": 36, "bold": True, "color": INKTXT, "font": "Cambria", "line": 1.0})
          for ln in name_lines])
    y += name_h + 0.08

    text(slide, MARGIN, y, col_w, 0.30, subtitle, size=13.5, color=BRICK)
    y += 0.36

    fob_w = min(col_w, 0.62 + 0.092 * len(fob))
    rect(slide, MARGIN, y, fob_w, 0.40, BRICK, rounded=True, adj=0.30)
    text(slide, MARGIN, y, fob_w, 0.40, fob, size=10.5, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.54

    about_h = 0.28 + 0.190 * len(about_lines) + 0.14
    rect(slide, MARGIN, y, col_w, about_h, INK, rounded=True, adj=0.045)
    text(slide, MARGIN + 0.28, y + 0.18, col_w - 0.56, 0.20, "ПРО КОМПАНІЮ",
         size=9.5, bold=True, color=GOLD, spacing=0.8)
    text(slide, MARGIN + 0.28, y + 0.43, col_w - 0.56, about_h - 0.5,
         [(ln, {"size": 10.2, "color": CREAM, "line": 1.16}) for ln in about_lines])
    y += about_h + 0.20

    n = len(stats)
    sw = col_w / n
    stat_top = y
    for i, (val, cap) in enumerate(stats):
        sx = MARGIN + i * sw
        if i:
            rect(slide, sx, stat_top + 0.03, 0.011, 0.60, LINE)
        pad = 0.20 if i else 0
        text(slide, sx + pad, stat_top, sw - pad - 0.1, 0.40, val, size=28, bold=True,
             color=BRICK, font="Cambria")
        text(slide, sx + pad, stat_top + 0.42, sw - pad - 0.1, 0.34, cap.split("\n"),
             size=8.3, color=GREY, line=1.14)
    y = stat_top + 0.80

    if flavours:
        y += 0.20
        cols = 3
        chip_gap = 0.08
        chip_w = (col_w - chip_gap * (cols - 1)) / cols
        chip_h = 0.32
        row_gap = 0.06
        for i, name in enumerate(flavours):
            c, r = i % cols, i // cols
            cx = MARGIN + c * (chip_w + chip_gap)
            cy = y + r * (chip_h + row_gap)
            oval(slide, cx + 0.02, cy + chip_h / 2 - 0.045, 0.09, FLAVOUR_DOT[name])
            text(slide, cx + 0.20, cy, chip_w - 0.22, chip_h, name, size=8.8, bold=True,
                 color=INKTXT, anchor=MSO_ANCHOR.MIDDLE)
        rows_used = -(-len(flavours) // cols)
        y += rows_used * (chip_h + row_gap)

    warn_h = 0.5
    warn_y = max(y + 0.22, 6.30)
    warn_y = min(warn_y, PAGE_H - 0.20 - warn_h)
    warning_bar(slide, warning, warn_y, h=warn_h, x=MARGIN, w=col_w)
    page_no(slide, page, on_dark_bar=(MARGIN + col_w - 0.62, warn_y, 0.47, warn_h))
    return slide


# ═══════════════════════════════════════════════════════════════════ SKU card
def sku(name, grams, per_carton, usd, eur, uah, desc, photo=None, photo_note=None):
    return dict(name=name, grams=grams, per_carton=per_carton, usd=usd, eur=eur,
                uah=uah, desc=desc, photo=photo, photo_note=photo_note)


def draw_card(slide, item, x, w, top):
    """Cards with a real photo sit on vivid brick; cards still waiting for one
    sit on near-black with a gold placeholder mark. Four identical red boxes
    in a row (the Takaokaya 45-60g slide) was the flattest part of v2 — this
    split turns "no photo yet" into a deliberate second colour, not a gap."""
    photo_h = 1.90
    pad = 0.16

    if item["photo"]:
        rect(slide, x, top, w, photo_h, BRICK)
        picture_contain(slide, item["photo"], x + 0.16, top + 0.14, w - 0.32, photo_h - 0.28)
        rect(slide, x + 0.10, top + 0.10, w - 0.20, photo_h - 0.20, None,
             line=GOLD, lw=1.0)
        if item["photo_note"]:
            rect(slide, x, top + photo_h - 0.22, w, 0.22, BRICK_D)
            text(slide, x + 0.08, top + photo_h - 0.22, w - 0.16, 0.22, item["photo_note"],
                 size=6.3, color=CREAM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    else:
        rect(slide, x, top, w, photo_h, INK_SOFT)
        placeholder_photo(slide, x, top, w, photo_h)

    inner_x, inner_w = x + pad, w - 2 * pad
    cy = top + photo_h + 0.20

    oval(slide, inner_x, cy + 0.035, 0.09, BRICK)
    text(slide, inner_x + 0.17, cy, inner_w - 0.17, 0.18,
         f"{num(item['grams'], 0)} Г · {item['per_carton']} ШТ/КАРТ.",
         size=7.3, bold=True, color=INKTXT, spacing=0.7)
    cy += 0.27

    per_line = max(12, int(inner_w / 0.098))
    lines = 1 + (len(item["name"]) - 1) // per_line
    name_h = 0.27 * lines
    text(slide, inner_x, cy, inner_w, name_h + 0.04, item["name"],
         size=14, bold=True, color=INKTXT, font="Calibri", line=1.03)
    cy += name_h + 0.09

    text(slide, inner_x, cy, inner_w, 0.62, item["desc"], size=8.3, color=GREY, line=1.16)

    price_y = top + photo_h + 1.62
    rect(slide, inner_x, price_y, inner_w, 0.86, INK, rounded=True, adj=0.06)
    text(slide, inner_x + 0.15, price_y + 0.10, inner_w - 0.30, 0.16,
         "ЗБІРНИЙ ВАНТАЖ 17 М³", size=6.3, bold=True, color=GOLD, spacing=0.6)
    text(slide, inner_x + 0.15, price_y + 0.29, inner_w * 0.56, 0.44,
         f"${num(item['usd'])}", size=21, bold=True, color=WHITE, font="Calibri")
    text(slide, inner_x + inner_w * 0.42, price_y + 0.42, inner_w * 0.58 - 0.15, 0.26,
         f"{num(item['uah'])} ₴", size=12.5, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    row_y = price_y + 0.98
    for label, value in (("С/С, €/од.", f"€{num(item['eur'])}"),
                         ("Питома собівартість", f"{num(item['uah'] / item['grams'])} ₴/г")):
        rect(slide, inner_x, row_y, inner_w, 0.010, LINE)
        text(slide, inner_x, row_y + 0.07, inner_w * 0.60, 0.20, label, size=7.3, color=GREY)
        text(slide, inner_x + inner_w * 0.40, row_y + 0.05, inner_w * 0.60, 0.22, value,
             size=9.3, bold=True, color=INKTXT, align=PP_ALIGN.RIGHT)
        row_y += 0.295

    return top + photo_h + 3.28   # card bottom, for reference


def card_slide(prs, layout, *, page, kicker_label, title, tag, items, warning):
    slide = blank_slide(prs, layout)
    kicker(slide, MARGIN, 0.34, kicker_label)
    text(slide, MARGIN, 0.58, 7.6, 0.60, title, size=32, bold=True, color=INKTXT, font="Cambria")

    tag_w = 3.55
    rect(slide, PAGE_W - MARGIN - tag_w, 0.52, tag_w, 0.46, BRICK, rounded=True, adj=0.30)
    text(slide, PAGE_W - MARGIN - tag_w, 0.52, tag_w, 0.46, tag, size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    n = len(items)
    top = 1.32
    gap = 0.22 if n == 4 else 0.30
    w = (PAGE_W - 2 * MARGIN - gap * (n - 1)) / n
    for i, item in enumerate(items):
        draw_card(slide, item, MARGIN + i * (w + gap), w, top)

    warn_y = 6.90
    warning_bar(slide, warning, warn_y, h=0.42, x=MARGIN, w=PAGE_W - 2 * MARGIN - 0.55)
    page_no(slide, page)
    return slide


# ═══════════════════════════════════════════════════════════════════ content
MS_SHRIMP, MS_NORI, MS_WASABI, TK_WASABI = (
    "assets/sku_ms_shrimp.png", "assets/sku_ms_nori.png",
    "assets/sku_ms_wasabi.png", "assets/sku_tk_wasabi.png")

MISHIMA_25 = [
    sku("Furikake з креветкою", 25, 20, 0.72, 0.8108, 36.49,
        "Видима креветка та дрібна риба — джерело мікроелементів. 3–4 г на 200 г рису.", MS_SHRIMP),
    sku("Furikake васабі", 25, 20, 0.89, 1.0022, 45.10,
        "Освіжаючий аромат васабі. Підходить до смаженого, м’яса на грилі та сашимі.",
        MS_WASABI, "фото — коробка 80 г з каталогу"),
    sku("Furikake зі смаком кімчі", 25, 20, 0.655, 0.7376, 33.19,
        "Гострий профіль кімчі. У роздрібному каталозі постачальника позиція відсутня.", None),
    sku("Furikake з норі", 25, 20, 0.555, 0.6250, 28.12,
        "Кунжут і норі — базовий смак серії. Універсальна присипка до рису.", MS_NORI),
]
MISHIMA_50 = [
    sku("Furikake з креветкою", 50, 64, 1.3797, 1.5537, 69.91,
        "Та сама рецептура у подвійній фасовці — 64 пакети в коробі.", MS_SHRIMP),
    sku("Furikake васабі", 50, 64, 1.6328, 1.8387, 82.74,
        "Найдорожча позиція Mishima в розрахунку: ₴1,65 за грам.",
        MS_WASABI, "фото — коробка 80 г з каталогу"),
    sku("Furikake зі смаком кімчі", 50, 64, 1.1969, 1.3478, 60.65,
        "Найнижча питома собівартість серед 50-грамових позицій.", None),
    sku("Furikake з норі", 50, 64, 0.9953, 1.1208, 50.44,
        "Найдешевша позиція розрахунку — ₴1,01 за грам.", MS_NORI),
]
TAKAOKAYA_BOTTLES = [
    sku("Nori Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
        "Класична норі-присипка у скляній пляшці fresh-lock.", None),
    sku("Nori Wasabi Furikake Bottle", 50, 10, 1.278, 1.9115, 86.02,
        "Норі з гранулами японського васабі (регіон Адзуміно). Ціна — як у поз. 1.",
        TK_WASABI, "фото роздрібної версії 70 г"),
    sku("Yuzu Kosho Furikake Bottle", 50, 10, 1.344, 2.0103, 90.46,
        "Юзу-кошьо з Кюсю: цитрусова цедра та зелений перець.", None),
]
TAKAOKAYA_45 = [
    sku("Nori Katsuo Furikake", 45, 80, 1.416, 2.1179, 95.31,
        "Норі та стружка боніто — умамі-профіль. Найвища питома с/с: ₴2,12/г.", None),
    sku("Nori Tamago Furikake", 60, 80, 1.416, 2.1179, 95.31,
        "Норі та яєчний порошок. Найбільша фасовка серії за тією ж ціною.", None),
    sku("Curry Furikake", 45, 80, 1.062, 1.5884, 71.48,
        "Каррі-присипка японського профілю спецій.", None),
    sku("Garlic Chili Oil Furikake", 45, 80, 1.062, 1.5884, 71.48,
        "Часник і чилі-олія раю — гострий напрямок серії. Ціна — як у поз. 6.", None),
]


# ═══════════════════════════════════════════════════════════════════ assemble
prs = Presentation(SRC)
blank = next(l for l in prs.slide_masters[0].slide_layouts if l.name == "Пустий слайд")

profile_slide(
    prs, blank, page=15, kicker_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 05",
    name_lines=["Dalian Mishima Foods", "Co., Ltd."],
    subtitle="Фурікаке, соуси та приправи Pan-Asian · Retail / OEM",
    fob="FOB DALIAN, КИТАЙ",
    photo="assets/mishima_profile_hero.jpg", photo_dark_bg=False, photo_caption=None,
    about_lines=[
        "Мішіма — виробник соусів, фурікаке, супових баз та юзу-продукції",
        "(м. Далянь, пров. Ляонін, КНР). Асортимент за каталогом: класичні",
        "соуси, пастові соуси, фурікаке в пакетах і коробках, супові бази,",
        "очадзуке, юзу-серія, фреш-лок пляшки.",
        "Контакт: foods@mishima.com.cn · +86 411 87611161",
    ],
    stats=[("8 SKU", "фурікаке в розрахунку\nсобівартості (25 і 50 г)"),
           ("17 м³", "єдиний прорахований\nсценарій доставки"),
           ("12 міс.", "термін придатності\n(за каталогом)")],
    flavours=None,
    warning="У файлі розрахунку базис доставки вказано «FOB Bangkok», хоча країна походження — Китай. "
            "Невідповідність потребує уточнення у постачальника до фіксації умов.",
)

card_slide(prs, blank, page=16, kicker_label="MISHIMA · FURIKAKE · 1/2",
          title="Furikake 25 г — собівартість", tag="ЗБІРНИЙ ВАНТАЖ 17 М³ · ОДИН СЦЕНАРІЙ",
          items=MISHIMA_25,
          warning="Розрахунок побудований лише для збірного вантажу 17 м³ — на відміну від "
                  "4-рівневого порівняння в профілі ZEK. Для зіставлення постачальників за обсягом "
                  "контейнера потрібно домоделювати сценарії 40′ / 20′ / 34 м³.")

card_slide(prs, blank, page=17, kicker_label="MISHIMA · FURIKAKE · 2/2",
          title="Furikake 50 г — собівартість", tag="ЗБІРНИЙ ВАНТАЖ 17 М³ · ОДИН СЦЕНАРІЙ",
          items=MISHIMA_50,
          warning="Смак «кімчі» відсутній у роздрібному каталозі постачальника — SKU та фото "
                  "потребують підтвердження. Фото 25 і 50 г ідентичні: постачальник показує одну "
                  "ілюстрацію пакування на смак.")

profile_slide(
    prs, blank, page=18, kicker_label="ПРОФІЛЬ ПОСТАЧАЛЬНИКА · 06",
    name_lines=["Takaokaya Co., Ltd.", "(бренд Kinjirushi)"],
    subtitle="Бутильоване фурікаке та васабі-приправи · Японія",
    fob="FOB ЯПОНІЯ (ПОРТ УТОЧНЮЄТЬСЯ)",
    photo="assets/takaokaya_profile_hero.jpg", photo_dark_bg=True,
    photo_caption="ВІЗУАЛІЗАЦІЯ ФАСОВКИ FRESH-LOCK · ФОТО ПОСТАЧАЛЬНИКА ВІДСУТНІ",
    about_lines=[
        "Заснована 1929 р., спеціалізується на продукції з автентичного",
        "японського васабі (регіон Адзуміно). Бренд Kinjirushi. Дочірня",
        "Takaokaya USA працює з 1986 р. Постачається у скляних «fresh-lock»",
        "пляшечках 45–60 г — преміум «Made in Japan» на відміну від",
        "пакетованого фурікаке з Китаю.",
    ],
    stats=[("1929", "рік заснування\nбренду"),
           ("45–60 г", "фасовка,\nскляна пляшка")],
    flavours=list(FLAVOUR_DOT),
    warning="Юридична особа й пряме контактне джерело постачальника в матеріалах компанії "
            "відсутні — профіль складено за відкритими даними бренду, потребує верифікації "
            "перед контрактом.",
)

card_slide(prs, blank, page=19, kicker_label="TAKAOKAYA · FURIKAKE · 1/2",
          title="Fresh-lock пляшки 50 г", tag="ЗБІРНИЙ ВАНТАЖ 17 М³ · ОДИН СЦЕНАРІЙ",
          items=TAKAOKAYA_BOTTLES,
          warning="Позиції 1 і 2 мають однакову закупівельну ціну незалежно від смаку — типово "
                  "для непроробленого прайсу; варто запросити поартикульну котирувку.")

card_slide(prs, blank, page=20, kicker_label="TAKAOKAYA · FURIKAKE · 2/2",
          title="Фасовка 45–60 г", tag="ЗБІРНИЙ ВАНТАЖ 17 М³ · ОДИН СЦЕНАРІЙ",
          items=TAKAOKAYA_45,
          warning="Пари 4/5 і 6/7 мають ідентичну ціну незалежно від ваги (45 г vs 60 г). Питома "
                  "собівартість Takaokaya (₴1,59–2,12/г) у 1,4–2,1 раза вища за Mishima "
                  "(₴0,99–1,80/г) — премія «Made in Japan».")

# The original deck had 7 slides: ZEK x3 (idx 0-2, kept) then the four
# Mishima/Takaokaya profile+table slides (idx 3-6, superseded). The 6 new
# slides were appended after them in creation order (idx 7-12), which is
# already the order we want, so dropping 3-6 needs no further reordering.
sld_lst = prs.slides._sldIdLst
ids = list(sld_lst)
for el in ids[3:7]:
    prs.part.drop_rel(el.rId)
    sld_lst.remove(el)

prs.save(DST)
print("saved", DST, "slides:", len(sld_lst))
